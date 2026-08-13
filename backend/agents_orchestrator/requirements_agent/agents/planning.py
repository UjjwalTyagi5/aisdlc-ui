from config.env import AGENTIC_BASE_URL
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import os
import time
import asyncio
import contextvars
import aiohttp
import aiofiles
import json
from bs4 import BeautifulSoup, NavigableString, Tag
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_COLOR_INDEX
from io import BytesIO
from xhtml2pdf import pisa
import textwrap
import docx 
import markdown
import litellm
from dotenv import load_dotenv
from typing import TypedDict, Annotated, List, Dict, Any, Optional
import PyPDF2
import shutil
from config.checkpoint import build_checkpointer as _build_checkpointer
from langchain_core.messages import BaseMessage, HumanMessage, ToolMessage, AIMessage
from langchain_core.tools import tool
from langchain_litellm import ChatLiteLLM
from langgraph.graph import StateGraph, END
from langgraph.graph.message import add_messages
from langgraph.prebuilt import ToolNode
from langgraph.prebuilt import InjectedState
from langgraph.types import Command
from shared.tools.mcp_runtime import get_mcp_tools, MCP_TOOLS_PROMPT_NOTE
from shared.services.skill_runtime import get_skill_tools
from config.ws_helper import set_session_id, broadcast_log, get_user_id, get_session_id
from config.connection_manager import manager
import os, sys, inspect
from agents_orchestrator.requirements_agent.prompts.brd_prompt import BRDPROMPT
from agents_orchestrator.requirements_agent.prompts.Pdd_prompt import PDDPROMPT
from agents_orchestrator.requirements_agent.prompts.MOM_prompt import MoMPROMPT
from agents_orchestrator.requirements_agent.prompts.risk_register_prompt import RISKPROMPT
from agents_orchestrator.requirements_agent.prompts.User_stories_prompt import USERSTORYPROMPT
from docxtpl import DocxTemplate, InlineImage, RichText
import pickle
from agents_orchestrator.requirements_agent.config import shared
from agents_orchestrator.requirements_agent.tools.nlp_quality_tool import run_nlp_quality_check
from agents_orchestrator.requirements_agent.tools.smell_tool import run_requirement_smell_check
from agents_orchestrator.requirements_agent.tools.spectral_tool import run_spectral_lint

import pandas as pd
from openpyxl.utils import get_column_letter

currentdir = os.path.dirname(os.path.abspath(inspect.getfile(inspect.currentframe())))
parentdir = os.path.dirname(currentdir)
superparentdir = os.path.dirname(parentdir)
supersuperparentdir = os.path.dirname(superparentdir)
from config import sdlcSettings
esett = sdlcSettings()

load_dotenv()


from shared.tools.document_tools import extract_file_text as _extract_file_text
from shared.errors import classify_error as _classify_error
from shared.services.artifact_service import write_and_notify as _write_and_notify
from shared.models.artifacts import RequirementsArtifact


# Hard cap on a single tool's output before it re-enters the LLM context. An
# uncapped result (a full board dump, a large file read) inflates every later
# turn's token cost and can silently blow the context window — which surfaces to
# the user as an opaque "the agent hit an error". Truncating keeps the loop
# stable; a tool that genuinely needs the full payload should persist it
# out-of-band and return a pointer, not stream megabytes back through the model.
_TOOL_OUTPUT_CAP = 12000


def _cap_tool_output(obs: Any) -> Any:
    """Truncate a tool observation to _TOOL_OUTPUT_CAP chars with a clear marker."""
    if not isinstance(obs, str):
        return obs
    if len(obs) <= _TOOL_OUTPUT_CAP:
        return obs
    dropped = len(obs) - _TOOL_OUTPUT_CAP
    return (
        obs[:_TOOL_OUTPUT_CAP]
        + f"\n\n…[output truncated — dropped {dropped} chars; capped at "
        f"{_TOOL_OUTPUT_CAP} to protect the context window]"
    )


@retry(
    retry=retry_if_exception_type(Exception),
    wait=wait_exponential(multiplier=1, min=2, max=30),
    stop=stop_after_attempt(3),
    reraise=True,
)
def _openai_generate(prompt: str, file_paths: list = None) -> str:
    """Backward-compatible generation helper. Uses the run's BYOK-resolved model
    (set on the resolved-model contextvar by the agent node). Fails closed if no
    model was resolved for this run."""
    from shared.services.model_resolver import get_resolved_model
    resolved = get_resolved_model()
    if resolved is None:
        raise RuntimeError(
            "No model resolved for this run — an administrator must configure a model provider.")
    parts = []
    for path in (file_paths or []):
        text = _extract_file_text(path)
        if text:
            parts.append(f"--- File: {os.path.basename(path)} ---\n{text}")
    content = "\n\n".join(parts) + "\n\n" + prompt if parts else prompt
    response = litellm.completion(
        model=resolved.model,
        custom_llm_provider=resolved.litellm_provider,
        api_key=resolved.api_key,
        api_base=resolved.base_url,
        messages=[{"role": "user", "content": content}],
        max_tokens=8192,  # client-facing BRD/PDD are long — 4096 truncated them mid-document
        temperature=0.15,
    )
    return response.choices[0].message.content or ""


async def _load_ref_paths(file_names: list) -> tuple:
    paths = []
    for name in file_names:
        if isinstance(name, str) and name.startswith("Error:"):
            continue
        if os.path.exists(name):
            try:
                async with aiofiles.open(name, "r") as f:
                    paths.append((await f.read()).strip())
            except Exception:
                continue
        else:
            return paths, f"Error: file {name} has not yet been uploaded"
    return paths, None


def add_hyperlink(paragraph, text, url):
    """
    Adds a hyperlink to a paragraph.
 
    Args:
        paragraph: The paragraph to add the hyperlink to.
        text (str): The text to display for the link.
        url (str): The URL the link should point to.
    """
    part = paragraph.part
    r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
   
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
   
    run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    rStyle = OxmlElement('w:rStyle')
    rStyle.set(qn('w:val'), 'Hyperlink')
    rPr.append(rStyle)
    run.append(rPr)
    run.append(OxmlElement('w:t', text=text))
   
    hyperlink.append(run)
    paragraph._p.append(hyperlink)
    return hyperlink

# Enhanced file_generated broadcast with file size
async def broadcast_file_generated(session_id: str, filename: str, file_path: str):
    """Enhanced function to broadcast file generation with file size"""
    try:
        user_id = get_user_id()
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        # Must match where the docx/pptx/png are actually written
        # ({FILES}/{user_id}/requirements_agent/{session_id}/output) — the static mount
        # is /generated → FILES_DIR. The old "/orchestrator/" segment 404'd every download.
        file_url = f"{AGENTIC_BASE_URL}/generated/{user_id}/requirements_agent/{session_id}/output/{filename}"
        await manager.broadcast({
            "type": "file_generated",
            "session_id": session_id,
            "filename": filename,
            "url": file_url,
            "file_size": file_size,
            "message": f"Generated file: {filename}"
        })

        print(f"DEBUG: Broadcasted file generation: {filename} ({file_size} bytes)")
        # Also persist as a project Artifact row so the file shows in the artifacts panel,
        # not only as a live in-chat download (best-effort; never blocks generation).
        try:
            from shared.services.chat_artifacts import register_generated_file  # noqa: PLC0415
            await register_generated_file(filename, file_path, file_url, stage="requirements")
        except Exception as _art_exc:  # noqa: BLE001
            print(f"DEBUG: register_generated_file failed: {_art_exc}")
        return file_url
    except Exception as e:
        print(f"ERROR: Failed to broadcast file generation: {str(e)}")
        return ""

async def markdown_to_docx(markdown_string: str, docx_path: str):
    """
    Converts a Markdown or HTML string to a DOCX file with near-full feature support.
    Now async to allow for real-time updates.
    """
    # Import OXML elements here to keep the function self-contained
    from docx.oxml.shared import OxmlElement
    from docx.oxml.ns import qn
 
    print(f"Starting markdown to docx conversion for: {docx_path}")
    dedented_markdown = textwrap.dedent(markdown_string).strip()
    doc = Document()
    # Enable all necessary extensions for full feature support
    html = markdown.markdown(dedented_markdown, extensions=['extra', 'sane_lists', 'tables', 'fenced_code', 'codehilite'])
    soup = BeautifulSoup(html, 'html.parser')
 
    def _add_inline_content(paragraph, element: Tag):
        """Recursively adds formatted runs to a paragraph."""
        for child in element.children:
            if isinstance(child, NavigableString):
                if child.strip():
                    run = paragraph.add_run(child)
                    curr = child.parent
                    # Apply formatting from parent tags
                    while curr is not None and curr.name != element.name:
                        if curr.name in ['strong', 'b']: run.bold = True
                        if curr.name in ['em', 'i']: run.italic = True
                        if curr.name in ['s', 'del']: run.strike = True
                        if curr.name == 'code':
                            run.font.name = 'Courier New'
                        curr = curr.parent
            elif isinstance(child, Tag):
                if child.name == 'br':
                    paragraph.add_run().add_break()
                elif child.name == 'a':
                    add_hyperlink(paragraph, child.get_text(), child.get('href'))
                else:
                    _add_inline_content(paragraph, child)

    async def _parse_element(element, doc, list_info=None):
        """Recursively parses a BeautifulSoup element and adds it to the document."""
        if not hasattr(element, 'name') or element.name is None: 
            return
 
        if element.name in [f'h{i}' for i in range(1, 7)]:
            level = int(element.name[1])
            p = doc.add_heading(level=level)
            _add_inline_content(p, element)
            print(f"Added heading level {level}")
            await asyncio.sleep(0.001)  # Allow other tasks to run
        elif element.name == 'p':
            p = doc.add_paragraph()
            _add_inline_content(p, element)
            print("Added paragraph")
            await asyncio.sleep(0.001)
        elif element.name in ['ul', 'ol']:
            list_style = 'List Bullet' if element.name == 'ul' else 'List Number'
            parent_level = list_info[1] if list_info else -1
            for li in element.find_all('li', recursive=False):
                await _parse_element(li, doc, list_info=(list_style, parent_level + 1))
        elif element.name == 'li':
            if list_info:
                style, level = list_info
                p = doc.add_paragraph(style=style)
                p.paragraph_format.left_indent = Inches(0.5 * level)
                _add_inline_content(p, element)
                for nested_list in element.find_all(['ul', 'ol'], recursive=False):
                    await _parse_element(nested_list, doc, list_info)
                print(f"Added list item at level {level}")
                await asyncio.sleep(0.001)
        elif element.name == 'table':
            rows_data = element.find_all('tr')
            if not rows_data: return
            num_cols = len(rows_data[0].find_all(['th', 'td']))
            table = doc.add_table(rows=len(rows_data), cols=num_cols)
            table.style = 'Table Grid'
            for i, row_element in enumerate(rows_data):
                cells = row_element.find_all(['th', 'td'])
                for j, cell_element in enumerate(cells):
                    p = table.cell(i, j).paragraphs[0]
                    _add_inline_content(p, cell_element)
            print(f"Added table with {len(rows_data)} rows and {num_cols} columns")
            await asyncio.sleep(0.001)
        elif element.name == 'img':
            src = element.get('src')
            if src:
                try:
                    async with aiohttp.ClientSession() as session:
                        async with session.get(src) as response:
                            if response.status == 200:
                                image_data = await response.read()
                                doc.add_picture(BytesIO(image_data), width=Inches(5.0))
                                print(f"Added image from {src}")
                            else:
                                print(f"Warning: Could not fetch image from {src}. Status: {response.status}")
                                doc.add_paragraph(f"[Image not found: {src}]").add_run().italic = True
                except Exception as e:
                    print(f"Warning: Could not fetch image from {src}. Error: {e}")
                    doc.add_paragraph(f"[Image not found: {src}]").add_run().italic = True
        elif element.name == 'hr':
            doc.add_page_break()
            print("Added page break")
            await asyncio.sleep(0.001)
        elif element.name == 'blockquote':
            p = doc.add_paragraph(style='Quote')
            p.paragraph_format.left_indent = Inches(0.5)
            for child in element.children:
                await _parse_element(child, doc)
            print("Added blockquote")
            await asyncio.sleep(0.001)
        elif element.name == 'pre':
            code_text = element.get_text()
            p = doc.add_paragraph(style='No Spacing')
            run = p.add_run(code_text)
            run.font.name = 'Courier New'
            run.font.size = Pt(10)
            print("Added code block")
            await asyncio.sleep(0.001)
 
    # Process all top-level tags from the parsed HTML
    for element in soup.contents:
        await _parse_element(element, doc)
    
    try:
        session_id = get_session_id()
        print(f"Saving document to {docx_path}")
        doc.save(docx_path)
        filename = os.path.basename(docx_path)
        _url = await broadcast_file_generated(session_id, filename, docx_path)
        print(f"Successfully saved document to {docx_path}")
        if _url:
            return f"Saved '{filename}'. Download it here: {_url}"
    except Exception as e:
        print(f"An exception occurred while saving the doc file: {e}")
    return f"Successfully converted Markdown to {os.path.basename(docx_path)}"

class AgentState(TypedDict):
    messages: Annotated[List[BaseMessage], add_messages]
    tenant_id: str
    model_id: str | None
    offering_id: str | None
    # The BYOK-resolved model for this run, carried through state so the sync `tools`
    # node (which runs in a fresh asyncio.run() context) can re-establish it for tools
    # that call _openai_generate — re-resolving there would hit the main-loop-bound DB.
    resolved_model: Any

@tool
async def upload_file(local_path: str):
    """
    Registers a local file so it can be used by other tools.

    Args:
        local_path (str): The local path of the file to register.

    Returns:
        str: A reference path that can be passed to other tools.
    """
    broadcast_log(manager, f"---Executing tool: Upload_file---", level="INFO")

    if not os.path.exists(local_path):
        return "Error: local file not found"

    try:
        user_id = get_user_id()
        session_id = get_session_id()
        input_dir = f"{esett.FILES}/{user_id}/requirements_agent/{session_id}/input"
        files_dir = os.path.join(input_dir, "files")
        os.makedirs(files_dir, exist_ok=True)

        filename = os.path.basename(local_path)
        dest_path = os.path.join(files_dir, filename)
        shutil.copy2(local_path, dest_path)

        ref_path = os.path.join(input_dir, filename + ".ref")
        async with aiofiles.open(ref_path, "w") as f:
            await f.write(dest_path)

        print(f"File registered: {dest_path}")
        return ref_path

    except Exception as e:
        error_msg = f"Exception during upload: {str(e)}"
        print(error_msg)
        return error_msg

@tool
async def delete_file(file_name: str):
    """
    Deletes a registered file reference and its local copy.

    Args:
        file_name (str): The reference path returned by upload_file.
    """
    print(f"---Executing tool: delete file for name: {file_name}---")
    try:
        if os.path.exists(file_name):
            async with aiofiles.open(file_name, "r") as f:
                actual_path = (await f.read()).strip()
            if os.path.exists(actual_path):
                os.remove(actual_path)
                print(f"Deleted file: {actual_path}")
            os.remove(file_name)
            print(f"Deleted reference: {file_name}")
            return f"Successfully deleted {os.path.basename(actual_path)}"
        else:
            return f"Warning: reference '{file_name}' not found, nothing deleted."
    except Exception as e:
        return f"Error deleting file: {e}"


# Per-session cache of the most recently generated document (BRD/PDD/Risk Register).
# The LLM reliably fails to echo a large generated markdown back as the markdowntodoc
# `content` arg (it calls it with {} — see empty-args validation error), so we stash the
# content here and let markdowntodoc fall back to it. Keyed by ws session id.
_LAST_GENERATED_DOC: dict[str, str] = {}


def _remember_doc(content: str) -> None:
    try:
        _LAST_GENERATED_DOC[get_session_id()] = content
    except Exception:
        pass


@tool
async def generate_brd(file_names: List[str], custom_prompt: str):
    """
    Generates a Business Requirement Document (BRD) by analyzing a collection of files,
    which can include transcripts, audios and meeting recordings.
 
    Args:
        file_names (List[str]): A list of file names to analyze
                                These names must be obtained from the upload_file tool
        custom_prompt (str): Additional instructions needed based on users query, can be empty
    """
    broadcast_log(manager, f"---Executing tool: generate_brd for files: {file_names}---", level="INFO")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    # Files are optional: generate from the user's description (custom_prompt) alone when
    # no file is uploaded. Only fail when there is nothing at all to work from.
    if not file_paths and not (custom_prompt and custom_prompt.strip()):
        return "Error: Provide a description of what the BRD should cover, or upload a source file."
    prompt = BRDPROMPT.format(custom_prompt=custom_prompt)
    try:
        broadcast_log(manager, "Generating BRD content...", level="INFO")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        broadcast_log(manager, "BRD generation completed", level="INFO")
        _remember_doc(response)
        return response
    except Exception as e:
        error_msg = _classify_error(e, "BRD generation")
        broadcast_log(manager, error_msg, level="ERROR")
        return error_msg

@tool
async def generate_pdd(file_names: List[str], custom_prompt: str):
    """
    Generates a Process Definition Document (PDD) by analyzing a collection of files,
    which can include transcripts, audios and meeting recordings.
   
    Args:
        file_names (List[str]): A list of file names to analyze
                                These names must be obtained from the upload_file tool
        custom_prompt (str): Additional special instructions to be taken into consideration while generating the PDD
    """
    print(f"Executing tool: generate_pdd for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths and not (custom_prompt and custom_prompt.strip()):
        return "Error: Provide a description of what the PDD should cover, or upload a source file."
    prompt = PDDPROMPT.format(custom_prompt=custom_prompt)
    try:
        print("Generating PDD content...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        print("PDD generation completed")
        _remember_doc(response)
        return response
    except Exception as e:
        return _classify_error(e, "PDD generation")

@tool
async def generate_risk_register(file_names: List[str], custom_prompt: str):
    """
    Generates a risk register by analyzing a collection of files,
    which can include transcripts, audios and meeting recordings.
   
    Args:
        file_names (List[str]): A list of file names to analyze
                                These names must be obtained from the upload_file tool
        custom_prompt (str): Additional special instructions to be taken into consideration while generating the PDD
    """
    print("Executing tool: generate_risk_register for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths and not (custom_prompt and custom_prompt.strip()):
        return "Error: Provide a description of what the Risk Register should cover, or upload a source file."
    prompt = RISKPROMPT.format(custom_prompt=custom_prompt)
    try:
        print("Generating Risk Register content...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        print("Risk Register generation completed")
        _remember_doc(response)
        return response
    except Exception as e:
        return _classify_error(e, "Risk Register generation")

@tool
async def template_pdd(content: str, files: List[str], filename : str = "template_pdd.docx"):
    """This is a function that takes a PDD related content and extracts, data from it to be filled into a predefined template
        This function also saves the document it generates
    Args:
        content (str) : The content from which the PDD template fields will be extracted, if file is provided can be empty
        files (List[str]) : A list of file names to analyze These names must be obtained from the upload_file tool
        filename (str) : The filename for saving, should be of the format project_topic_template_pdd.docx
    """
    file_paths, err = await _load_ref_paths(files)
    if err:
        return err

    prompt = f"""
You are a specialist data extraction AI. Your task is to read the unstructured text provided, comprehend its meaning, and populate a precise JSON schema with the information you synthesize.
 
**CRITICAL DIRECTIVES:**
1.  **ADHERE TO THE SCHEMA:** You MUST use the exact field names as defined in the JSON schema in Part 2. Do not change, add, or omit keys.
2.  **UNSTRUCTURED INPUT:** The source text is free-form (prose, a transcript, etc.). It will likely NOT contain tables. You must infer the data from sentences and paragraphs.
3.  **JSON OUTPUT ONLY:** Your entire response MUST be a single, valid JSON object, without any surrounding text or markdown.
4.  **HANDLE MISSING DATA:** If information for any field or object is not mentioned, use an empty string `""` or an empty list `[]` as appropriate. Do not invent data.
 
---
**PART 1: CONCEPTUAL ANALYSIS GUIDE**
---
*This section explains how to find the information needed for the specific keys in the schema.*
 
*   **For lists (`risks`, `assumptions`, `functional_req`, etc.):** Scan the document for mentions of each item. For every distinct item you find (e.g., a specific risk), create a JSON object.
    *   **`"no"`**: If the text numbers the item (e.g., "1. The first risk is..."), use that number. Otherwise, assign your own sequential number (as a string: "1", "2", "3"...).
    *   **`"desc"`**: This is the main description of the item. For a risk, this is the description of the risk itself. For a functional requirement, it's the requirement's description.
    *   **`"impact"`, `"priority"`, `"dependency"`, `"remarks"`**: Find any text that describes these specific attributes for the item and place it in the corresponding field.
*   **For document tables (`invoice`, `seal`, `booking`, `shipping`):** The text will describe these fields in prose. For example, "The invoice's field for the shipper should have a header of 'Exporter'." From this, you would create an object: `{{"no": "1", "field": "shipper", "header": "Exporter", "remarks": ""}}`.
*   **For metadata (`project_title`, `author`, etc.):** Identify these single pieces of information from the text.
 
---
**PART 2: FINAL JSON OUTPUT SCHEMA (STRICT)**
---
*Your entire output must be a single JSON object that strictly follows this schema.*
 
```json
{{{{
  "project_title": "(string)",
  "author": "(string)",
  "version": "(string)",
  "date": "(string)",
  "version_table": [
    {{{{
      "version": "(string)",
      "date": "(string)",
      "changes": "(string)",
      "author": "(string)",
      "reviewer": "(string)",
      "approver": "(string)",
      "signoff": "(string)"
    }}}}
  ],
  "introduction": "(string)",
  "intended_audience": "(string)",
  "objective": "(string)",
  "solution_objective": "(string)",
  "functional_table": [
    {{{{
      "business": "(string)",
      "description": "(string)"
    }}}}
  ],
  "steps": [
    {{{{
      "no": "(string)",
      "action": "(string)",
      "details": "(string)"
    }}}}
  ],
  "bl_document": [
    {{{{
      "no": "(string)",
      "document": "(string)",
      "file": "(string)"
    }}}}
  ],
  "invoice": [
    {{{{
      "no": "(string)",
      "field": "(string)",
      "header": "(string)",
      "remarks": "(string)"
    }}}}
  ],
  "seal": [
    {{{{
      "no": "(string)",
      "field": "(string)",
      "header": "(string)",
      "remarks": "(string)"
    }}}}
  ],
  "booking": [
    {{{{
      "no": "(string)",
      "field": "(string)",
      "header": "(string)",
      "remarks": "(string)"
    }}}}
  ],
  "shipping": [
    {{{{
      "no": "(string)",
      "field": "(string)",
      "header": "(string)",
      "remarks": "(string)"
    }}}}
  ],
  "business_logic": [
    {{{{
      "no": "(string)",
      "field": "(string)",
      "business_rule": "(string)",
      "remarks": "(string)"
    }}}}
  ],
  "success_criteria": [
    {{{{
      "criteria": "(string)",
      "target": "(string)",
      "description": "(string)"
    }}}}
  ],
  "success_scope": "(string)",
  "functional_req": [
    {{{{
      "no": "(string)",
      "desc": "(string)",
      "priority": "(string)",
      "dependency": "(string)"
    }}}}
  ],
  "out_of_scope": [
    {{{{
      "no": "(string)",
      "desc": "(string)",
      "priority": "(string)",
      "dependency": "(string)"
    }}}}
  ],
  "assumptions": [
    {{{{
      "no": "(string)",
      "desc": "(string)",
      "impact": "(string)"
    }}}}
  ],
  "risks": [
    {{{{
      "no": "(string)",
      "desc": "(string)",
      "impact": "(string)"
    }}}}
  ],
  "constraints": [
    {{{{
      "no": "(string)",
      "desc": "(string)"
    }}}}
  ],
  "dependencies": [
    {{{{
      "no": "(string)",
      "desc": "(string)",
      "impact": "(string)"
    }}}}
  ]
}}}}"""
    
    print("Extracting data for PDD template...")
    loop = asyncio.get_event_loop()
    user_id = get_user_id()
    session_id = get_session_id()
    full_prompt = prompt + "\n\nAdditional context:\n" + content
    ctx = contextvars.copy_context()
    response = await loop.run_in_executor(None, ctx.run, _openai_generate, full_prompt, file_paths)

    s = response.replace("```json", "")
    s = s.replace("```", "")
    s = json.loads(s)
    
    print("Filling PDD template...")
    file = "PDD_Template.docx"
    doc = DocxTemplate(file)
    doc.render(s)
    word_doc = doc.docx
    out_path = f"{esett.FILES}/{user_id}/requirements_agent/{session_id}/output"
    if os.path.exists(out_path):
            print(f"The folder '{out_path}' already exists.")
    else:
        os.makedirs(out_path)
    output_path = os.path.join(out_path, filename)
    
    await loop.run_in_executor(None, doc.save, output_path)
    print(f"PDD template saved to {output_path}")
    try:
        session_id = get_session_id()
        print(f"Saving document to {output_path}")
        doc.save(output_path)
        filename = os.path.basename(output_path)
        await broadcast_file_generated(session_id, filename, output_path)
        print(f"Successfully saved document to {output_path}")

    except Exception as e:
        print(f"An exception occurred while saving the doc file: {e}")
    return f"Document saved to {output_path}"

@tool
async def update_response(file_names: List[str], query: str, content: str):
    """
    A function that can be used to update the content provided based on context, returns the updated content
 
    Args:
    file_names (List[str]): A list of file names obtained from upload file that is provided as context
    query (str): the actual query or the updates the user needs, combine all the needs previously mentioned
    content (str): the actual content that needs to be updated
    """
    print(f"Executing tool: update_response for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    prompt = f"""Given the following content and attached files as context update the content provided to you based on the query
Query: {query}

````````
Content:
{content}
`````````
    """
    try:
        print("Updating response based on query...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        print("Response update completed")
        return response
    except Exception as e:
        return _classify_error(e, "document update")

@tool
async def generate_mom(file_names: List[str], custom_prompt: str):
    """
    Generates the Minutes of the meeting by analyzing a collection of files,
    which can include transcripts, audios and meeting recordings.
   
    Args:
        file_names (List[str]): A list of file names to analyze
                                These names must be obtained from the upload_file tool
        custom_prompt (str): Additional special instructions to be taken into consideration while generating the PDD
    """
    print("Executing tool: generate_mom for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths:
        return "Error: Must provide at least one file to generate the minutes of the meeting"
    prompt = MoMPROMPT.format(custom_prompt=custom_prompt)
    try:
        print("Generating Minutes of Meeting...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        print("MOM generation completed")
        return response
    except Exception as e:
        return _classify_error(e, "MOM generation")

@tool
async def general_query(file_names: List[str], query: str):
    """
    Can be used to generate a response for a general user query. A general user query is one
    that does ask to generate a specific type of document rather is just looking for an answer
   
    Args:
        file_names (List[str]): A list of file names to analyze
                                These names must be obtained from the upload_file tool
        query (str): the users query
    """
    print(f"Executing tool: general_query for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths:
        return "Error: Must provide at least one file to answer the query"
    try:
        print("Processing general query...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, query, file_paths)
        print("General query processing completed")
        return response
    except Exception as e:
        return _classify_error(e, "general query")
    
@tool
async def generate_user_stories(file_names: List[str], custom_prompt: str = ""):
    """
    Generates user stories by analyzing a collection of files,
    which can include BRDs, meeting notes, requirements documents, or customer feedback.
    Args:
    file_names (List[str]): A list of file names to analyze
    These names must be obtained from the upload_file tool
    custom_prompt (str): Additional special instructions for user story generation
    """
    broadcast_log(manager, f"---Executing tool: generate_user_stories for files: {file_names}---", level="INFO")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths:
        return "Error: Must provide at least one valid file to generate user stories"
    prompt = USERSTORYPROMPT.format(custom_prompt=custom_prompt)
    try:
        broadcast_log(manager, "Generating user stories content...", level="INFO")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        text = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, file_paths)
        broadcast_log(manager, "User stories generation completed", level="INFO")
        try:
            if text.strip().startswith('[') or text.strip().startswith('{'):
                return json.dumps(json.loads(text), indent=2)
            return text
        except json.JSONDecodeError:
            return text
    except Exception as e:
        error_msg = _classify_error(e, "user story generation")
        broadcast_log(manager, error_msg, level="ERROR")
        return error_msg

@tool
async def revise_user_stories(existing_stories: str, feedback: str, custom_prompt: str = ""):
    """
    Revises existing user stories based on user feedback and suggestions.
    Args:
        existing_stories (str): The current user stories in JSON or text format
        feedback (str): User feedback describing what changes to make
        custom_prompt (str): Additional special instructions for the revision
    """
    print(f"---Executing tool: revise_user_stories with feedback: {feedback[:100]}...---")
    revision_context = f"""
    CURRENT USER STORIES:
    {existing_stories}
    USER FEEDBACK FOR REVISION:
    {feedback}
    Please revise the user stories based on the feedback provided.
    """
    prompt = USERSTORYPROMPT.format(custom_prompt=custom_prompt + " " + revision_context)
    try:
        print("Revising user stories based on feedback...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        text = await loop.run_in_executor(None, ctx.run, _openai_generate, prompt, [])
        print("User stories revision completed")
        try:
            if text.strip().startswith('[') or text.strip().startswith('{'):
                return json.dumps(json.loads(text), indent=2)
            return text
        except json.JSONDecodeError:
            return text
    except Exception as e:
        return _classify_error(e, "user story revision")

@tool
async def generate_planning_sheet(file_names: List[str], custom_prompt: str = "", filename: str = "planning_sheet.xlsx"):
    """
    Generates a planning sheet in Excel format by extracting data from a collection of files.
    The model will only extract existing data and will not invent any information.
    The planning sheet will have columns such as Milestone, Task, Owner, start, end, and planned duration.
    This Function automatically saves the planning sheet
    Args:
        file_names (List[str]): A list of file names to analyze.
                                These names must be obtained from the upload_file tool.
        custom_prompt (str): Additional instructions for the model, can be empty.
        filename (str): filename for saving, should be of the format project_topic_planning_sheet.xlsx
    """
    print(f"---Executing tool: generate_planning_sheet for files: {file_names}---")
    file_paths, err = await _load_ref_paths(file_names)
    if err:
        return err
    if not file_paths:
        return "Error: Must provide at least one file to generate a planning sheet."

    planning_prompt = f"""
    Analyze the content of the provided files and extract information to create a planning sheet.
    The planning sheet should have the following columns: Milestone, Task, Owner, start, end, planned duration.
    You must only extract information that is explicitly mentioned in the files. Do not invent or infer any data.
    If information for a specific field is not available in the documents, use words like TBD or None.
    Return the extracted data in a JSON format where each key is a column name and the value is a list of the data for that column.
    Only output the json do not include any other words in your response.
    {custom_prompt}
    """

    try:
        print("Generating planning sheet data...")
        loop = asyncio.get_event_loop()
        ctx = contextvars.copy_context()
        response = await loop.run_in_executor(None, ctx.run, _openai_generate, planning_prompt, file_paths)
        s = response
        s = s.replace("```json", "")
        s = s.replace("```", "")
        print("Planning sheet data generated, creating Excel file...")
        planning_data = pd.read_json(s)
            
        loop = asyncio.get_event_loop()
        user_id = get_user_id()
        session_id = get_session_id()
        out_path = f"{esett.FILES}/{user_id}/requirements_agent/{session_id}/output"
        if os.path.exists(out_path):
            print(f"The folder '{out_path}' already exists.")
        else:
            os.makedirs(out_path)
    
        output_path = os.path.join(out_path, filename)
        await loop.run_in_executor(None, _save_excel_file, planning_data, output_path)
        try:
            session_id = get_session_id()
            print(f"Saving sheet to {output_path}")
            filename = os.path.basename(output_path)
            await broadcast_file_generated(session_id, filename, output_path)
            print(f"Successfully saved sheet to {output_path}")

        except Exception as e:
            print(f"An exception occurred while saving the sheet file: {e}")
        print(f"Planning sheet saved to {output_path}")
        return f"Successfully generated planning sheet as {filename}"
 
    except Exception as e:
        return _classify_error(e, "planning sheet generation")

def _save_excel_file(planning_data, output_path):
    """Helper function to save Excel file in executor"""
    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        planning_data.to_excel(writer, index=False, sheet_name="PlanningSheet")
        worksheet = writer.sheets["PlanningSheet"]
        
        for col_idx, col in enumerate(planning_data.columns, 1):
            max_length = max(planning_data[col].astype(str).map(len).max(), len(col))
            worksheet.column_dimensions[get_column_letter(col_idx)].width = max_length + 2
 
@tool
async def markdowntodoc(content: str = "", output_path: str = "requirements_document.docx"):
    """
    Convert an HTML or Markdown string into a docx file and save it (downloadable in chat).
    Args:
    content (str): the markdown string to convert. If omitted or empty, the most recently
        generated document (BRD / PDD / Risk Register) for this session is used automatically —
        so you do NOT need to paste the whole document back in.
    output_path (str): output filename, e.g. project_topic_filename.docx
    """
    # Fall back to the last generated doc when the model calls this with no/empty content.
    if not (content and content.strip()):
        try:
            content = _LAST_GENERATED_DOC.get(get_session_id(), "")
        except Exception:
            content = ""
    if not (content and content.strip()):
        return ("Error: nothing to convert. Generate a BRD / PDD / Risk Register first, "
                "or pass the markdown content explicitly.")
    if not (output_path and output_path.strip()):
        output_path = "requirements_document.docx"

    print(output_path,"@@@@@@@#111111111111")
    output_path = os.path.basename(output_path)
    user_id = get_user_id()
    session_id = get_session_id()
    out_path = f"{esett.FILES}/{user_id}/requirements_agent/{session_id}/output"
    print(out_path,"@@@@@@@@@@@@@@111111222222222222")
    if os.path.exists(out_path):
        print(f"The folder '{out_path}' already exists.")
    else:
        os.makedirs(out_path)
    output_path = os.path.join(out_path, output_path)
    print(output_path,"-------------------")

    return await markdown_to_docx(content, output_path)


def _requirements_output_dir() -> str:
    """{FILES}/{user}/requirements_agent/{session}/output — where generated files that
    are downloadable from chat must live (matches broadcast_file_generated + the
    /generated static mount)."""
    out = f"{esett.FILES}/{get_user_id()}/requirements_agent/{get_session_id()}/output"
    os.makedirs(out, exist_ok=True)
    return out


def _download_bytes(url: str) -> Optional[bytes]:
    """Fetch bytes from a URL (used to pull a rendered diagram PNG)."""
    try:
        import httpx as _httpx  # noqa: PLC0415
        r = _httpx.get(url, timeout=30, follow_redirects=True)
        r.raise_for_status()
        return r.content
    except Exception:
        return None


@tool
async def generate_ppt(content: str, output_path: str) -> str:
    """Generate a PowerPoint (.pptx) deck and make it downloadable from the chat.

    Structure the content in markdown: each '#'/'##' heading starts a new slide and
    '- ' lines become bullet points. The first slide is used as the title slide.

    Args:
        content: Markdown-style deck content (headings = slides, bullets = points).
        output_path: Filename for the deck, e.g. 'procurement_overview.pptx'.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        return "Error: PowerPoint generation needs python-pptx — run `pip install python-pptx` on the backend."

    session_id = get_session_id()
    filename = os.path.basename(output_path) or "presentation.pptx"
    if not filename.lower().endswith(".pptx"):
        filename += ".pptx"
    dest = os.path.join(_requirements_output_dir(), filename)

    def _build() -> None:
        prs = Presentation()
        # Parse markdown into (title, [bullets]) slides.
        slides: List[tuple] = []
        cur_title: Optional[str] = None
        cur_bullets: List[str] = []
        for raw in content.splitlines():
            s = raw.strip()
            if not s:
                continue
            if s.startswith("#"):
                if cur_title is not None:
                    slides.append((cur_title, cur_bullets))
                cur_title, cur_bullets = s.lstrip("#").strip(), []
            elif s.startswith(("- ", "* ")):
                cur_bullets.append(s[2:].strip())
            else:
                cur_bullets.append(s)
        if cur_title is not None:
            slides.append((cur_title, cur_bullets))
        if not slides:
            slides = [("Presentation", [content.strip()[:400]])]

        first_title, first_bullets = slides[0]
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = first_title
        if first_bullets and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = first_bullets[0]

        for title, bullets in slides[1:]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets or [""]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = b
        prs.save(dest)

    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, _build)
    url = await broadcast_file_generated(session_id, filename, dest)
    tail = f" Download it here: {url}" if url else ""
    return f"Created PowerPoint deck '{filename}'.{tail}"


@tool
async def generate_diagram(diagram_type: str, source: str, output_path: str = "") -> str:
    """Render a flowchart/diagram as a PNG image (via kroki.io) and make it downloadable
    from the chat — flowcharts, sequence, ER, class, or architecture sketches.

    Args:
        diagram_type: e.g. mermaid, plantuml, graphviz, c4plantuml, erd, bpmn, excalidraw.
        source: The diagram source code in that syntax.
        output_path: Optional filename, e.g. 'process_flow.png'.
    """
    import base64 as _b64  # noqa: PLC0415
    import zlib as _zlib  # noqa: PLC0415

    dtype = diagram_type.strip().lower()
    compressed = _b64.urlsafe_b64encode(_zlib.compress(source.encode("utf-8"), 9)).decode("utf-8")
    png_url = f"https://kroki.io/{dtype}/png/{compressed}"

    loop = asyncio.get_event_loop()
    img = await loop.run_in_executor(None, _download_bytes, png_url)
    if not img:
        return f"Error: could not render the {dtype} diagram — check the diagram type/syntax."

    session_id = get_session_id()
    filename = os.path.basename(output_path) or f"diagram_{dtype}.png"
    if not filename.lower().endswith(".png"):
        filename += ".png"
    dest = os.path.join(_requirements_output_dir(), filename)
    with open(dest, "wb") as fh:
        fh.write(img)
    url = await broadcast_file_generated(session_id, filename, dest)
    tail = f" Download it here: {url}" if url else ""
    return f"Rendered {dtype} diagram '{filename}'.{tail}"

# Convert tools to async
# tools = [upload_file, delete_file, generate_brd, generate_mom, generate_pdd, 
#          generate_risk_register, general_query, update_response, markdowntodoc, 
#          template_pdd, generate_planning_sheet]

@tool
async def write_requirements_artifact(
    run_id: str,
    brd_content: str = "",
    user_stories: list = None,
    acceptance_criteria: list = None,
    risk_register: list = None,
) -> str:
    """Persist the requirements artifact to the database and notify the orchestrator.

    Call this when the user confirms requirements are complete to advance the
    pipeline to the design stage.

    Args:
        run_id: Pipeline run identifier (provided in the system context).
        brd_content: The full BRD markdown content.
        user_stories: List of user story dicts.
        acceptance_criteria: List of acceptance criteria strings.
        risk_register: List of risk register dicts.
    """
    session_id = get_session_id()
    artifact = RequirementsArtifact(
        agent_session_id=session_id or run_id,
        brd_content=brd_content or None,
        user_stories=user_stories,
        acceptance_criteria=acceptance_criteria,
        risk_register=risk_register,
    )
    await _write_and_notify(run_id, "requirements", artifact.model_dump())
    return f"Requirements artifact persisted for run {run_id}"


# ── Board connector tools (ADO/Jira ingestion) ──────────────────────────────
# These let the agent pull epics/stories from the tenant's connected board. The
# connector is injected per-run via config.connectors.context.set_connector()
# (stage runner / Redis worker); these tools read it back with get_connector().
# Credentials never appear here — the connector resolves the tenant PAT ephemerally
# inside auth_adapter(). When no board is connected the tools fail closed with a
# clear, user-facing message instead of raising (so a run without a board still
# proceeds on uploaded/pasted input).
from config.connectors.context import get_connector as _get_active_connector


async def _board_connector():
    """Return (connector, None) using the run-injected connector, or (None, error_str)."""
    try:
        connector = _get_active_connector()
    except Exception:
        return None, (
            "No project-management board is connected for this run. An administrator "
            "must connect Azure DevOps or Jira on the Integrations page before board "
            "stories can be ingested. You can still proceed with pasted or uploaded "
            "requirements."
        )
    try:
        broadcast_log(
            manager,
            f"Using {connector.display_name} board connector ({connector.connector_name})",
            level="INFO",
        )
    except Exception:
        pass
    return connector, None


@tool
async def list_board_projects() -> str:
    """List all projects visible in the connected project-management board.

    Call this first when the user wants to ingest stories from the board.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        projects = await connector.read_adapter("list_projects")
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching projects: {exc}"
    if not projects:
        return "No projects found."
    lines = [
        f"- {p.get('name', '')}" + (f" ({p.get('key')})" if p.get("key") else "")
        for p in projects
    ]
    return f"Projects available in {connector.display_name}:\n" + "\n".join(lines)


@tool
async def list_board_groups(project: str) -> str:
    """List optional board groups/teams for the given project.

    Jira usually returns no groups (its project roles are permissions, not delivery
    teams). If none are returned, continue with list_board_items(project).

    Args:
        project: Exact project name/key as returned by list_board_projects.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        teams = await connector.read_adapter("list_teams", project=project)
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching teams for project '{project}': {exc}"
    if not teams:
        return (
            f"No board groups/teams found in project '{project}'. "
            "Continue without a team filter by calling list_board_items(project)."
        )
    lines = [f"- {t['name']}" for t in teams]
    return f"Teams in '{project}':\n" + "\n".join(lines)


@tool
async def list_board_states(project: str, work_item_type: str = "User Story") -> str:
    """List the workflow states defined for a work item type in the project.

    Use this to discover real state names (e.g. "New", "Active") instead of guessing.
    Defaults to User Story.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        states = await connector.read_adapter(
            "list_states", project=project, item_type=work_item_type
        )
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching states: {exc}"
    if not states:
        return "No states returned."
    lines = [f"- {s['name']} ({s.get('category', '')})" for s in states]
    return f"States for '{work_item_type}' in '{project}':\n" + "\n".join(lines)


@tool
async def list_board_items(project: str, team: Optional[str] = None) -> str:
    """Fetch ALL work items in a project regardless of state or type.

    Use this FIRST when the user asks to pull existing board items — it never
    misses items due to unknown state names or unsupported types.

    Args:
        project: Project name.
        team: Optional team name to scope the query.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        items = await connector.read_adapter("list_all_items", project=project, team=team)
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching work items: {exc}"
    if not items:
        return f"No work items found in project '{project}'."
    lines = [
        f"#{s['source_key']} (id: {s['id']}) [{s['work_item_type']}] [{s['state']}] {s['title']}"
        + (f" (assigned: {s['assigned_to']})" if s.get("assigned_to") else "")
        for s in items
    ]
    return f"Found {len(items)} work items in '{project}':\n" + "\n".join(lines)


@tool
async def list_board_items_by_state(project: str, state: str, team: Optional[str] = None) -> str:
    """List user stories matching a state filter in the project.

    Args:
        project: Project name.
        state: Exact state name from list_board_states (e.g. "New", "Active", "To Do").
        team: Optional team name to scope the query.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        stories = await connector.read_adapter(
            "list_stories", project=project, state=state, team=team
        )
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching stories: {exc}"
    if not stories:
        return f"No stories in state '{state}' for project '{project}'."
    lines = [
        f"#{s['source_key']} (id: {s['id']}) [{s['work_item_type']}] {s['title']}"
        + (f" (assigned: {s['assigned_to']})" if s.get("assigned_to") else "")
        for s in stories
    ]
    return f"Found {len(stories)} stories in '{state}':\n" + "\n".join(lines)


@tool
async def fetch_board_item_detail(project: str, work_item_id: int) -> str:
    """Fetch the full normalized detail of one story (description + acceptance criteria)."""
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        norm = await connector.read_adapter(
            "fetch_item_detail", project=project, item_id=work_item_id
        )
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching work item #{work_item_id}: {exc}"
    ac_lines = "\n".join(f"  - {c}" for c in norm.get("acceptance_criteria", [])) or "  (none)"
    return (
        f"Story #{norm.get('source_key') or norm['work_item_id']} (id: {norm['id']}): {norm['title']}\n"
        f"State: {norm['state']} | Type: {norm['work_item_type']}\n"
        f"Description:\n{norm.get('description') or '(empty)'}\n"
        f"Acceptance Criteria:\n{ac_lines}"
    )


@tool
async def fetch_board_hierarchy(project: str) -> str:
    """Fetch the full Epic → Feature → User Story hierarchy from the board.

    Use when the user wants all work items organised by parent-child relationships
    rather than a flat list.

    Args:
        project: Exact project name.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        tree = await connector.read_adapter("fetch_hierarchy", project=project)
    except Exception as exc:  # noqa: BLE001
        return f"Error fetching hierarchy: {exc}"
    if not tree:
        return f"No work items found in project '{project}'."

    def _render(nodes: list, indent: int = 0) -> list:
        out = []
        prefix = "  " * indent
        for node in nodes:
            ac_count = len(node.get("acceptance_criteria", []))
            out.append(
                f"{prefix}[{node['type']}] #{node['id']} {node['title']} "
                f"({node['state']}, {ac_count} AC)"
            )
            if node.get("children"):
                out.extend(_render(node["children"], indent + 1))
        return out

    return "\n".join([f"Work item hierarchy for '{project}':"] + _render(tree))


@tool
async def create_board_item(
    project: str,
    title: str,
    work_item_type: str = "User Story",
    description: str = "",
    acceptance_criteria: str = "",
) -> str:
    """Create a single work item of ANY type (Epic/Feature/User Story/Task/Bug) on the board.

    Args:
        project: Project name.
        title: Work item title.
        work_item_type: Exact type — e.g. "Epic", "Feature", "User Story", "Task", "Bug".
        description: Optional description.
        acceptance_criteria: Optional acceptance criteria text.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        wi = await connector.write_adapter(
            "create_item",
            project=project,
            item_type=work_item_type,
            title=title,
            description=description,
            acceptance_criteria=acceptance_criteria,
        )
        wid = wi.get("work_item_id") or wi.get("id", "?")
        url = wi.get("url") or wi.get("work_item_url") or wi.get("_links", {}).get("html", {}).get("href", "")
        return f"Created {work_item_type} #{wid}: {title}" + (f"\n{url}" if url else "")
    except Exception as exc:  # noqa: BLE001
        return f"Error creating {work_item_type} '{title}': {exc}"


@tool
async def create_board_project(
    name: str,
    key: str = "",
    description: str = "",
) -> str:
    """Create a NEW project/space on the connected board (Jira project or ADO project).

    Requires an admin-scoped connector token. On Jira a project key is auto-derived from
    the name when not given; on ADO the project provisions asynchronously.

    Args:
        name: The new project's display name.
        key: Optional Jira project key (uppercase, ≤10 chars); ignored for ADO.
        description: Optional project description.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        res = await connector.write_adapter(
            "create_project", name=name, key=key, description=description
        )
        if res.get("queued"):
            return f"Project '{name}' creation queued on Azure DevOps (provisions in the background)."
        return f"Created project '{res.get('name', name)}' (key: {res.get('key', '?')})."
    except Exception as exc:  # noqa: BLE001
        return f"Error creating project '{name}': {exc}"


@tool
async def move_board_item_state(project: str, work_item_ids: List[int], target_state: str) -> str:
    """Move one or more work items to a workflow state. Confirm with the user first.

    Args:
        project: Project name.
        work_item_ids: List of work item IDs to move.
        target_state: Target state name (use list_board_states for valid values).
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    moved, failed = [], []
    for wid in work_item_ids:
        try:
            await connector.write_adapter(
                "move_item_state", project=project, item_id=wid, new_state=target_state
            )
            moved.append(f"#{wid}")
        except Exception as exc:  # noqa: BLE001
            failed.append(f"#{wid}: {exc}")
    lines = [f"Moved {len(moved)} item(s) to '{target_state}':"] + moved
    if failed:
        lines += [f"Failed ({len(failed)}):"] + failed
    return "\n".join(lines)


@tool
async def add_board_comment(project: str, work_item_id: int, comment: str) -> str:
    """Add a comment to a work item — for gap reports, design links, PR URLs, or audit notes.

    Args:
        project: Project name.
        work_item_id: Work item ID to comment on.
        comment: Comment text to post.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        await connector.write_adapter(
            "add_comment", project=project, item_id=work_item_id, comment=comment
        )
        return f"Comment added to #{work_item_id}."
    except Exception as exc:  # noqa: BLE001
        return f"Error adding comment to #{work_item_id}: {exc}"


@tool
async def update_board_item(
    project: str,
    work_item_id: str,
    title: str = "",
    description: str = "",
    acceptance_criteria: str = "",
    state: str = "",
) -> str:
    """Update an existing work item — change its title, description, acceptance criteria,
    and/or move its workflow state. Only the fields you pass are changed.

    Args:
        project: Project name (Jira: the project name or key, e.g. "My Software Team").
        work_item_id: Work item id (ADO, e.g. "42") or Jira issue key (e.g. "SCRUM-5").
        title: New title/summary (optional).
        description: New description (optional).
        acceptance_criteria: New acceptance criteria (optional; applied on ADO).
        state: New workflow state to transition to (optional; use list_board_states for valid values).
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    changed: List[str] = []
    try:
        if title or description or acceptance_criteria:
            await connector.write_adapter(
                "update_item",
                project=project,
                item_id=work_item_id,
                title=title,
                description=description,
                acceptance_criteria=acceptance_criteria,
            )
            changed.append("fields")
        if state:
            await connector.write_adapter(
                "move_item_state", project=project, item_id=work_item_id, new_state=state
            )
            changed.append(f"state→{state}")
        if not changed:
            return f"Nothing to update for #{work_item_id} — pass a title, description, acceptance_criteria, or state."
        return f"Updated #{work_item_id}: {', '.join(changed)}."
    except Exception as exc:  # noqa: BLE001
        return f"Error updating #{work_item_id}: {exc}"


@tool
async def delete_board_item(project: str, work_item_id: str) -> str:
    """Permanently delete a work item from the board. IRREVERSIBLE — confirm the id with
    the user first. (ADO: moved to the project Recycle Bin; Jira: the issue is deleted.)

    Args:
        project: Project name (Jira: the project name or key).
        work_item_id: Work item id (ADO, e.g. "42") or Jira issue key (e.g. "SCRUM-5").
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        await connector.write_adapter("delete_item", project=project, item_id=work_item_id)
        return f"Deleted work item #{work_item_id}."
    except Exception as exc:  # noqa: BLE001
        return f"Error deleting #{work_item_id}: {exc}"


# ── §4.1 authoring tools (gap analysis · normalisation · epics · write-back) ────
# These implement the AGENT_LIFECYCLE_REFERENCE §4.1 capability set conversationally
# on top of the live board connector. LLM tools run the BYOK-resolved model off the
# event loop; write-back tools go through the same provider abstraction as the board
# tools above (create_item / update_item_fields / add_comment).


@tool
async def normalize_acceptance_criteria(stories_json: str) -> str:
    """Rewrite raw acceptance criteria from board stories into Gherkin Given/When/Then.

    Call after fetch_board_item_detail / fetch_board_hierarchy (or with the
    selected_stories from context) to turn free-form AC into structured, testable
    scenarios. Generates 2-4 scenarios for stories that have none.

    Args:
        stories_json: JSON array of stories — each with "id", "title", and
            "acceptance_criteria" (list of raw strings).

    Returns a JSON array with the same stories but acceptance_criteria replaced
    with normalised Gherkin scenarios.
    """
    prompt = f"""You are a senior QA engineer. Produce complete, testable acceptance criteria
in Gherkin format (Given/When/Then) for each user story.

INPUT STORIES (JSON):
{stories_json}

Rules:
1. If a story already has acceptance_criteria items: rewrite each into one Gherkin scenario
   ("Given ... When ... Then ..."). Clean up vague phrasing — keep the original intent.
2. If a story has an empty acceptance_criteria array: GENERATE 2-4 concrete Gherkin scenarios
   based on the story title and description. Cover the happy path and at least one edge case.
   Do NOT leave any story with an empty acceptance_criteria list.
3. Add an "And" clause when there are multiple outcomes in the same scenario.
4. Do NOT invent requirements beyond what the story title/description implies.
5. Each scenario must be independently testable by a QA engineer.

Return a JSON array. Each element:
{{"id": <original work item id>, "title": "<story title>",
  "description": "<story description or empty string>",
  "acceptance_criteria": ["Given ... When ... Then ...", ...]}}

No markdown. Only valid JSON."""
    loop = asyncio.get_event_loop()
    ctx = contextvars.copy_context()
    return await loop.run_in_executor(None, ctx.run, _openai_generate, prompt)


@tool
async def detect_requirement_gaps(stories_json: str) -> str:
    """Analyse stories and detect gaps, ambiguities, and missing NFRs before design.

    Use after fetching/normalising stories to ensure requirements are complete before
    handing off to the Design Agent. Returns a completeness score and ONLY genuine
    follow-up questions (max 3).

    Args:
        stories_json: JSON array of stories with title, description, and
            acceptance_criteria fields.

    Returns a structured gap report (JSON) with actionable follow-up questions.
    """
    prompt = f"""You are a senior business analyst reviewing user stories before system design begins.

CRITICAL RULES — read before analysing:
1. Only flag a gap if the information is GENUINELY absent from the story text, description, and acceptance criteria.
   If the story already states the answer, do NOT raise a question about it.
2. Do NOT ask about NFRs (performance, security, accessibility, scalability, availability) unless the story
   explicitly involves a novel performance requirement or a new security boundary.
3. Do NOT ask about scope explicitly ruled out in the story.
4. Maximum 3 gaps. If you cannot find 3 genuine missing items, return fewer. Do not pad with generic questions.
5. A gap must be something a developer or designer cannot proceed without.

STORIES:
{stories_json}

Return a JSON object:
{{"completeness_score": <0-100>, "ready_for_design": <true|false>,
  "gaps": [{{"story_id": <id or null>, "story_title": "<title or 'General'>",
    "gap_type": "<missing_ac|ambiguous|scope_gap|conflict>",
    "description": "<what specific information is absent>",
    "follow_up_question": "<exact question to ask>"}}],
  "summary": "<one paragraph overall assessment>"}}

No markdown. Only valid JSON."""
    loop = asyncio.get_event_loop()
    ctx = contextvars.copy_context()
    return await loop.run_in_executor(None, ctx.run, _openai_generate, prompt)


@tool
async def identify_epics(stories_text: str) -> str:
    """Group user stories into Epics based on theme and business domain.

    Args:
        stories_text: User stories text or JSON.

    Returns a JSON array of epics, each with epic_title, epic_description, and stories.
    """
    prompt = f"""You are a product manager. Group the following user stories into logical Epics.

STORIES:
{stories_text}

Output a JSON array. Each element:
- "epic_title": short name for the epic
- "epic_description": one sentence describing the business goal
- "stories": list of story titles that belong to this epic

No markdown, no extra text. Only valid JSON."""
    loop = asyncio.get_event_loop()
    ctx = contextvars.copy_context()
    return await loop.run_in_executor(None, ctx.run, _openai_generate, prompt)


@tool
async def generate_stories_from_brd(brd_content: str, project_context: str = "") -> str:
    """Generate INVEST user stories + Gherkin AC directly from BRD / requirements text.

    Args:
        brd_content: Full BRD / requirements document text.
        project_context: Optional project name, tech stack, or constraints.

    Returns JSON with "epics" (grouped stories) and "stories" (flat INVEST list with AC).
    """
    prompt = f"""You are a senior business analyst. Analyse the BRD below and extract all functional requirements.
For each requirement produce a user story and acceptance criteria.

BRD:
{brd_content}

{"PROJECT CONTEXT: " + project_context if project_context else ""}

Return a single JSON object:
{{"epics": [{{"epic_title": "...", "epic_description": "...", "stories": ["story title 1"]}}],
  "stories": [{{"title": "As a <role>, I want <action> so that <benefit>", "description": "...",
    "epic": "parent epic title", "acceptance_criteria": ["Given ... When ... Then ..."],
    "story_points_estimate": <fibonacci number>,
    "invest_check": {{"Independent": true, "Negotiable": true, "Valuable": true, "Estimable": true, "Small": true, "Testable": true}}}}]}}
No markdown. Only valid JSON."""
    loop = asyncio.get_event_loop()
    ctx = contextvars.copy_context()
    return await loop.run_in_executor(None, ctx.run, _openai_generate, prompt)


@tool
async def write_stories_to_board(stories_json: str, project: str) -> str:
    """Create generated user stories as NEW work items on the board (bulk create).

    Args:
        stories_json: JSON array of stories (from generate_stories_from_brd /
            generate_user_stories). Each story needs at least a "title"; optional
            "description" and "acceptance_criteria" (list).
        project: Board project name to create the stories in.

    Returns a summary of created work item IDs.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        stories = json.loads(stories_json)
        if isinstance(stories, dict) and "stories" in stories:
            stories = stories["stories"]
    except Exception as exc:  # noqa: BLE001
        return f"Error parsing stories JSON: {exc}"

    created, failed = [], []
    for s in stories:
        title = (s.get("title") or "").strip()
        if not title:
            continue
        ac_lines = s.get("acceptance_criteria", [])
        ac_html = "<br>".join(ac_lines) if isinstance(ac_lines, list) else str(ac_lines)
        try:
            wi = await connector.write_adapter(
                "create_item",
                project=project,
                item_type="User Story",
                title=title,
                description=s.get("description", ""),
                acceptance_criteria=ac_html,
            )
            created.append(f"#{wi['id']} {title}")
        except Exception as exc:  # noqa: BLE001
            failed.append(f"{title}: {exc}")

    lines = [f"Created {len(created)} work item(s):"] + created
    if failed:
        lines += [f"\nFailed ({len(failed)}):"] + failed
    return "\n".join(lines)


@tool
async def write_acceptance_criteria_to_board(
    ac_json: str, story_ids: List[int], project: str
) -> str:
    """Write generated acceptance criteria back onto EXISTING work items.

    Args:
        ac_json: JSON array of {story_title, acceptance_criteria} entries.
        story_ids: Ordered list of work item IDs matching the entries in ac_json.
        project: Board project name.

    Returns a summary of updated work item IDs.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        ac_list = json.loads(ac_json)
    except Exception as exc:  # noqa: BLE001
        return f"Error parsing AC JSON: {exc}"

    updated, failed = [], []
    for idx, wid in enumerate(story_ids):
        if idx >= len(ac_list):
            break
        criteria = ac_list[idx].get("acceptance_criteria", [])
        ac_html = "<br>".join(criteria) if isinstance(criteria, list) else str(criteria)
        try:
            await connector.write_adapter(
                "update_item_fields", project=project, item_id=wid, acceptance_criteria=ac_html
            )
            updated.append(f"#{wid}")
        except Exception as exc:  # noqa: BLE001
            failed.append(f"#{wid}: {exc}")

    lines = [f"Updated {len(updated)} work item(s):"] + updated
    if failed:
        lines += [f"\nFailed ({len(failed)}):"] + failed
    return "\n".join(lines)


@tool
async def write_back_normalized_to_board(
    stories_json: str, project: str, add_audit_comment: bool = True
) -> str:
    """Write normalised descriptions + Gherkin AC back to existing work items.

    Each story must have an "id" (work item ID) plus updated "description" and/or
    "acceptance_criteria". Optionally adds an audit comment noting the AI normalisation.

    Args:
        stories_json: JSON array of stories from normalize_acceptance_criteria.
        project: Board project name.
        add_audit_comment: If True, adds an audit comment to each updated work item.

    Returns a summary of updated work item IDs.
    """
    connector, err = await _board_connector()
    if err:
        return f"Error: {err}"
    try:
        stories = json.loads(stories_json)
        if isinstance(stories, dict) and "stories" in stories:
            stories = stories["stories"]
    except Exception as exc:  # noqa: BLE001
        return f"Error parsing stories JSON: {exc}"

    from datetime import datetime as _dt
    audit_note = (
        f"Requirements normalised by SDLC AI Agent on {_dt.utcnow().strftime('%Y-%m-%d %H:%M')} UTC. "
        "Acceptance criteria rewritten to Gherkin (Given/When/Then) format."
    )
    updated, failed = [], []
    for s in stories:
        wid = s.get("id")
        if not wid:
            continue
        ac_list = s.get("acceptance_criteria", [])
        ac_html = "<br>".join(ac_list) if isinstance(ac_list, list) else str(ac_list)
        try:
            await connector.write_adapter(
                "update_item_fields",
                project=project,
                item_id=wid,
                acceptance_criteria=ac_html,
                description=s.get("description", ""),
            )
            if add_audit_comment:
                try:
                    await connector.write_adapter(
                        "add_comment", project=project, item_id=wid, comment=audit_note
                    )
                except Exception:  # noqa: BLE001 — comment is best-effort
                    pass
            updated.append(f"#{wid} {s.get('title', '')}")
        except Exception as exc:  # noqa: BLE001
            failed.append(f"#{wid}: {exc}")

    lines = [f"Updated {len(updated)} work item(s):"] + updated
    if failed:
        lines += [f"\nFailed ({len(failed)}):"] + failed
    return "\n".join(lines)


@tool
async def build_requirements_payload(
    project: str,
    team: str,
    stories_json: str,
    scope_summary: str = "",
    assumptions: str = "",
    out_of_scope: str = "",
) -> str:
    """Build the structured requirements payload the Design Agent will consume.

    Call this as the FINAL packaging step (after normalisation + gap resolution) and
    before persisting via write_requirements_artifact. Produces a machine-readable
    summary of all gathered requirements.

    Args:
        project: Board project name/key.
        team: Optional board team/group name.
        stories_json: JSON array of stories with normalised AC.
        scope_summary: 1-2 sentence description of what is being built.
        assumptions: Comma-separated assumptions (or empty string).
        out_of_scope: Comma-separated out-of-scope items (or empty string).

    Returns a JSON payload string prefixed with REQUIREMENTS_PAYLOAD::.
    """
    from datetime import datetime as _dt

    try:
        stories = json.loads(stories_json)
        if isinstance(stories, dict) and "stories" in stories:
            stories = stories["stories"]
    except Exception:  # noqa: BLE001
        stories = []

    kind = "unknown"
    connector, err = await _board_connector()
    if not err and connector is not None:
        try:
            kind = connector.connector_name
        except Exception:  # noqa: BLE001
            kind = "unknown"

    work_items = []
    for s in stories:
        wid = s.get("id") or s.get("work_item_id")
        if wid:
            work_items.append({
                "id": int(wid),
                "source_key": s.get("source_key") or s.get("key") or str(wid),
                "provider_kind": kind,
                "type": s.get("work_item_type", "User Story"),
                "title": s.get("title", ""),
            })

    payload = {
        "generated_at": _dt.utcnow().isoformat() + "Z",
        "project": project,
        "team": team,
        "provider_kind": kind,
        "scope_summary": scope_summary or f"Requirements from project '{project}'",
        "story_count": len(stories),
        "stories": stories,
        "work_items": work_items,
        "assumptions": [a.strip() for a in assumptions.split(",") if a.strip()],
        "out_of_scope": [o.strip() for o in out_of_scope.split(",") if o.strip()],
        "non_functional_requirements": {},
    }
    return f"REQUIREMENTS_PAYLOAD::\n{json.dumps(payload, indent=2)}"


_BOARD_TOOLS = [
    list_board_projects, list_board_groups, list_board_states, list_board_items,
    list_board_items_by_state, fetch_board_item_detail, fetch_board_hierarchy,
    create_board_project,
    create_board_item, update_board_item, delete_board_item,
    move_board_item_state, add_board_comment,
    normalize_acceptance_criteria, detect_requirement_gaps, identify_epics,
    generate_stories_from_brd, write_stories_to_board, write_acceptance_criteria_to_board,
    write_back_normalized_to_board, build_requirements_payload,
]


# Convert tools to async
tools = [upload_file, delete_file, generate_brd, generate_mom, generate_pdd,
         generate_risk_register, general_query, update_response, markdowntodoc,
         generate_ppt, generate_diagram,
         template_pdd, generate_planning_sheet, generate_user_stories, revise_user_stories,
         write_requirements_artifact,
         run_nlp_quality_check, run_requirement_smell_check, run_spectral_lint,
         *_BOARD_TOOLS]


# ── System prompt ──────────────────────────────────────────────────────────────

INGESTION_SYS_MESSAGE = """\
You are the Requirements & Ingestion Agent. Your sole responsibility is to
gather, normalise, validate, and package project requirements — then hand off
to the Design Agent when the user is ready.

── CAPABILITIES ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
A. PM INGESTION        — pull stories, features, and epics from {PM_PROVIDER}
B. HIERARCHY VIEW      — show full Epic → Feature → User Story tree
C. STORY GENERATION    — generate INVEST-formatted user stories from any input
D. AC NORMALISATION    — rewrite raw AC into Gherkin Given/When/Then format
E. GAP DETECTION       — identify missing AC, NFRs, and ambiguities
E2. NLP PRE-VALIDATION  — run deterministic quality checks BEFORE drafting
F. REQUIREMENTS PAYLOAD — build structured JSON for the Design Agent
G. EPIC IDENTIFICATION  — group stories into Epics by business theme
H. BRD → STORIES       — analyse a pasted BRD and produce full story breakdown
I. WRITE TO PM TOOL     — create/update work items and write back normalised content

── NLP PRE-VALIDATION (run BEFORE story drafting from raw input) ──────────────────
Before you draft user stories or a BRD from raw input, run the deterministic
pre-validation tools on the source text and surface their findings to the user:
  • run_nlp_quality_check  — flags vague/weak terms, passive voice, entities
  • run_requirement_smell_check — flags untestable verbs, compound & ambiguous requirements
  • run_spectral_lint — lint any referenced OpenAPI spec file (skip if none / unavailable)
Present flagged items as open questions or assumptions to confirm with the user. These tools
are advisory and degrade gracefully — if one reports "unavailable", proceed with the rest.

── ACCURACY RULES (READ FIRST — CRITICAL FOR DEMO) ──────────────────────────────────────────────────────────────────────────────────────────
These rules prevent hallucination. Violating them produces wrong output.

1. ONLY OUTPUT WHAT IS IN THE SOURCE MATERIAL.
   - User stories, AC, business rules, personas: derive ONLY from what the user
     typed, pasted, uploaded, or what the project management tool returned. Never invent them.
   - If a detail (e.g. a rule, a threshold, a role) is not in the source,
     surface it to the user as an open question or a stated assumption —
     do NOT silently fill it in.

2. GAP REPORT = HONEST GAPS, NOT INVENTED ANSWERS.
   - detect_requirement_gaps (when the user explicitly requests it) must report
     actual missing information.
   - Never pre-populate a gap answer yourself. Ask the user. Record their answer.
   - If a gap analysis was run, only call build_requirements_payload AFTER the
     user has answered every follow-up question from it. If no gap analysis was
     requested, package directly — list unresolved unknowns as assumptions.

3. CONFIRM INFERENCES BEFORE COMMITTING.
   - If you infer something not explicitly stated (e.g. "I assume this is a
     web app"), state the inference and ask: "Is that correct?"
   - Do not proceed until the user confirms or corrects it.

4. GENERATED STORIES MUST TRACE TO INPUT.
   - Every generated user story must cite its source: the BRD section, the
     user's message, or the work item ID it came from.
   - Format: "Source: [Board item #1234 or KEY-123 / BRD Section 3.2 / User provided]"

── SCOPE BOUNDARY (CRITICAL) ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────
You are responsible for requirements. You do NOT produce design or code yourself.
- If the user asks to proceed to development or build the code: package requirements
  and emit HANDOFF to development. Never block this request.
- If the user asks about HLD, LLD, API design, C4 diagrams, or wireframes:
  say "That's handled by the Design Agent" and emit HANDOFF to design.
- Never produce design documents or code yourself.

── SELECTED-STORY SCOPE (READ — drives which work items you touch) ─────────────────────────────
The structured pipeline context may include a `requirements` object. When present it
carries the stories the user has ALREADY PULLED from their board, so you DO NOT need to
re-pull or ask which project:
  • `board_project` — the ADO project already chosen/finalised by the user. Use THIS
    for every board operation. Do NOT ask the user which project, and do NOT call
    list_board_projects, when `board_project` is set.
  • `selected_stories` — the full content (ref, title, description, acceptance_criteria)
    of the stories the user selected. `selected_story_refs` is the same list of refs.
  • `all_stories` — ref + title of every pulled story; `all_story_refs` the refs.

Rules:
- If `selected_stories` / `selected_story_refs` is present AND non-empty: you ALREADY
  HAVE these stories' content — do NOT fetch them again and do NOT ask which project.
  Scope ALL of your work (analyse, normalise, gap-check, edit, comment, write-back) to
  ONLY those items. Begin by acknowledging the specific story you're working on (by title).
- If nothing is selected but `all_stories` is non-empty: operate on all pulled
  stories (you already have their refs/titles; fetch detail via fetch_board_item_detail
  using `board_project` only when you need a description you weren't given).
- ONLY ask the user which project / to pull stories when the `requirements` object is
  absent or has NO stories at all (board_project empty AND all_stories empty) — i.e.
  nothing has been ingested yet.
- Create new stories on ADO via create_board_item in `board_project` — never just
  describe them. After creating, confirm the created item ids/titles in your reply.
  Never tell the user to click a "Pull stories" button — that UI does not exist here.
- Edit existing selected stories on ADO (add_board_comment for notes; create_board_item
  for children). ADO is the source of truth: never invent a story not written back to ADO.

── FLOW A: PM INGESTION (Standard) ─────────────────────────────────────────────────────────────────────────────────────────────────────────────
1. Call list_board_projects. If one project, use it; otherwise ask user.
   Project output may include both display name and key, e.g. "SDLC Team (SCRUM)".
   Use the exact display name or key returned by the tool.
2. If the provider is Jira, DO NOT require or invent a team. Jira project roles
   like Member, Viewer, Administrator, or atlassian-addons-project-access are
   permissions, not delivery teams. Skip team filtering unless the user gave a
   real board team field explicitly.
3. For Azure DevOps only, call list_board_groups for the project. If one group/team exists, use it; otherwise ask.
4. Call list_board_items(project, team) — fetches ALL items regardless
   of state or type. NEVER guess states. This is the correct first call.
   Present results as a numbered list with ID, type, state, and title.
5. Ask which stories to process (user can say "all", "user stories only", or pick IDs).
6. Run FLOW G — start with the PRE-STEP to fetch full details, then normalise + gap check.
   Do NOT skip the PRE-STEP even if the user says "process all" — you still need to
   fetch each story's description and AC before normalising.
7. When the user is satisfied, package the work: call build_requirements_payload with
   the normalised stories, then persist it with write_requirements_artifact (passing
   run_id, brd_content if generated, user_stories, acceptance_criteria, risk_register).
   Offer to write the normalised stories back to the board (write_back_normalized_to_board)
   and to move their state (move_board_item_state) — always confirm first.

RULE: Do NOT call list_board_states + list_board_items_by_state to find initial board items.
      Always use list_board_items — it never misses items due to
      unknown state names or unsupported work item types.
RULE: Never print bracketed fake tool narration like "[Called list_board_items]".
      Only summarize results that came from an actual tool response.
RULE: Never write simulated user replies such as "User: Process all stories",
      "User: Member", or any other fake next message. End with a direct question
      to the real user and wait.

── FLOW B: HIERARCHY VIEW ────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
- When the user says "show hierarchy", "show epics", "show all work items":
  Call fetch_board_hierarchy(project) to display the provider's parent-child tree.
- Present the tree clearly. Offer to fetch full details for any node.

── FLOW G: NORMALISE + PACKAGE (run after any story fetch) ────────────────────────────────────────────────────────────────────────────

FAST-PATH FOR EXPLICIT DEVELOPMENT INTENT:
If the current user message already says to start/proceed with development,
skip all review prompts and blocking gap-question stops. Fetch full story
details if needed, normalize acceptance criteria, call build_requirements_payload,
then emit HANDOFF to development immediately. Do not ask "does this look
correct?", do not ask gap questions, and do not ask for another confirmation.
You may include assumptions/gaps as non-blocking risks inside the payload/summary.

PRE-STEP — FETCH FULL STORY DETAILS (mandatory):
Before normalising, you need description + existing AC for each story.
- The flat list from list_board_items contains only titles — it has no description or AC.
- For each User Story (NOT Epics or Features) in the selected set:
  Call fetch_board_item_detail(project, work_item_id) to retrieve the full content.
- Build a JSON array from the results: [{"id": N, "title": "...", "description": "...", "acceptance_criteria": [...]}, ...]
  Use this enriched JSON for all subsequent steps. Never run normalize on title-only data.

1. Call normalize_acceptance_criteria on the enriched stories JSON.
   The tool generates Gherkin AC from the story title/description when no AC exists yet.

   Present the results as formatted markdown — NOT raw JSON. Use this structure for each story:

   ### Story #[id]: [title]
   **Acceptance Criteria (Gherkin):**
   - Given [context] When [action] Then [outcome]
   - Given [context] When [action] Then [outcome]

   *** STOP HERE. Ask: "Does this look correct? Any changes needed?" ***
   Wait for the user's answer before proceeding.

2. Once the user confirms the stories/AC, assemble the REQUIREMENTS DOCUMENT and present it
   as formatted markdown — scope summary, the stories with their Gherkin AC, assumptions,
   and out-of-scope items. Then ask: "Requirements look complete — proceed to Design, or
   skip straight to Development?"  *** STOP HERE — wait for the user's routing answer only. ***

3. As soon as the user answers (e.g. "design" / "development" / "yes, development"):
   call build_requirements_payload immediately, then emit the matching HANDOFF.
   Do NOT ask for a second confirmation. See HANDOFF RULES for the exact sentinel.
4. Offer to write the normalised AC back to the board provider (write_back_normalized_to_board) BEFORE
   emitting the HANDOFF if the user has not already declined it, but do not block the
   HANDOFF on this — if they don't respond or say skip, emit HANDOFF immediately.
5. Emit the HANDOFF signal as the very last line of your response.

GAP ANALYSIS — ON EXPLICIT REQUEST ONLY:
- NEVER call detect_requirement_gaps unless the user explicitly asks for it
  (e.g. "check for gaps", "run gap analysis", "anything missing?").
- Do NOT offer or announce gap analysis as the next step of the flow.
- When the user DOES ask, run detect_requirement_gaps on the enriched stories JSON and present:

   **Completeness Score: [X]/100**
   **Gaps found: [N]**
   1. [Story #id] [question]

   *** STOP HERE. Do NOT call any more tools. Do NOT answer your own questions.
       Your response must end with the questions. Wait for the user's answers. ***
- After the user answers each question, incorporate their answer.
  If the user says "it's already in the story", "that's defined", "already covered",
  or any equivalent — accept it immediately, mark that gap closed, and move to the next.
  Do NOT re-ask or ask for clarification on a closed gap.
  If there are more unanswered questions, ask the next one and STOP again.
  NEVER bundle all questions into one message. One question at a time.
- When all gaps are resolved, return to step 2 of FLOW G (requirements document + routing question).

── FLOW C/D: GENERATE STORIES, AC & EPICS ─────────────────────────────────────────────────────────────────────────────────────────────────
- If user provides story titles/details: call generate_user_stories.
- Call normalize_acceptance_criteria on the result.
- Call identify_epics to group them.
- Present ALL generated stories with their Gherkin AC to the user.
  *** STOP HERE. Ask the user if the stories look correct before proceeding. ***
- Ask if user wants to push to the board provider via write_stories_to_board.
- Then run FLOW G steps 2–5 (requirements document → routing → payload → HANDOFF).
- Gap analysis only if the user explicitly asks (see FLOW G gap rules).

CRITICAL — SELF-ANSWERING IS FORBIDDEN:
- NEVER answer your own follow-up questions.
- NEVER pre-fill gap answers based on assumptions.
- If you asked a clarifying question, your response MUST END with that question.
  Do not write anything after it. Do not call any more tools until the user replies.

CRITICAL — BULK GAP DISMISSAL:
- If the user says anything like "everything is in the story", "all info is already there",
  "that's final", "no more questions needed", or "it's all defined" — treat ALL remaining
  gap questions as closed simultaneously.
  Do NOT ask any more gap questions. Jump straight to the finalised summary and HANDOFF prompt.

CRITICAL — GAP QUALITY:
- The gap report may contain up to 3 genuine gaps maximum.
- Only ask about something if a developer literally cannot build it without the answer.
- NFRs covered by existing system standards are NOT gaps.

── FLOW H: BRD → STORIES ────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
- Ask user to paste the BRD text (or key sections).
- Call generate_stories_from_brd with the pasted content.
- Run normalize_acceptance_criteria on the result.
- Present epics, then stories with Gherkin AC.
- Ask if user wants to create these as board items.
- Gap analysis only if the user explicitly asks (see FLOW G gap rules).

── FLOW I: WRITE TO PM TOOL ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
- create_board_item: creates ONE board item of ANY type.
  If the user says "create an epic", use work_item_type="Epic".
  If they say "create a feature", use work_item_type="Feature".
- write_stories_to_board: bulk-creates User Stories from generated JSON.
- write_back_normalized_to_board: updates existing board items with normalised
  description + Gherkin AC and adds an audit comment.
- write_acceptance_criteria_to_board: patches AC onto existing board item IDs.
- move_board_item_state: move one or more items to a specific state.
  Use after processing starts (→ "Active") and after user confirms (→ "Ready for Design").
- add_board_comment: post any text as a comment on a board item.
  Use to attach gap report summaries to the parent epic.
- ALWAYS confirm with the user before calling any mutating tool.

── FLOW J: GENERATE DOCUMENTS ───────────────────────────────────────────────────────────────────────────────────────────────────────────────
After requirements are confirmed and the gap report is reviewed, offer to export:
- generate_brd_document(project, scope_summary, stories_json, gap_report_json)
  → full BRD as DOCX, returns download URL
- generate_user_stories_document(project, stories_json)
  → all stories with Gherkin AC as DOCX, returns download URL
- generate_risk_register_document(project, stories_json, gap_report_json)
  → risk register based on stories and gap analysis, returns download URL
Always present the download URL from the tool result as a clickable link in your response.
Post the gap report summary as a comment on the parent epic/item using add_board_comment.

── FILE UPLOADS ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
- When the user's message says "please use the following files ..." followed by
  one or more file paths, call read_uploaded_file for EACH path first.
- Pass extracted text to generate_stories_from_brd or generate_user_stories.
- Supported formats: .txt, .md, .csv, .xlsx, .xls, .docx, .pdf.
- NEVER say you cannot access files — always call read_uploaded_file.

── HANDOFF RULES (CRITICAL) ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
ALWAYS call build_requirements_payload before emitting any HANDOFF.

ROUTING — determine destination from the user's intent:

  → "to":"development"  when user says anything like: "proceed with development",
    "skip design", "just build it", "go to dev", "start coding", "yes" (after
    being asked Design-or-Development and they chose Development), or any message
    expressing intent to go directly to coding. DO NOT block this — respect the
    user's choice immediately.

  → "to":"design"  when user says: "proceed to design", "create design",
    "architecture", "yes" (after being asked Design-or-Development and they chose
    Design), or any message expressing intent to create the design first.

The sequence before any HANDOFF:
  a) Call build_requirements_payload (if not already called this session).
  b) Show a brief summary of finalised requirements.
  c) Ask ONE question: "Requirements look complete — proceed to Design, or skip
     straight to Development?"  This is the user's ONLY confirmation step.
  d) As soon as the user answers — emit the matching HANDOFF immediately.
     Do NOT ask "are you sure?", do NOT ask them to confirm again, do NOT say
     "shall I proceed?" — the answer to step c IS the confirmation.
     The frontend Approve button provides a second safety gate; you must not
     also ask for a typed confirmation.

SHORTCUT — if the user's message ALREADY expresses routing intent before step c
  (e.g. "proceed with development", "go to dev", "skip design"), skip step c
  entirely and emit HANDOFF immediately after step b.

HANDOFF signals:

1. To Design (user wants design first):
   HANDOFF::{"to":"design","batch_id":"DESIGN-READY","stage_completed":"requirements","context_keys":["requirements_payload"],"triggered_by":"user_confirmed"}

2. To Development (user explicitly skips design):
   HANDOFF::{"to":"development","batch_id":"DEV-READY","stage_completed":"requirements","context_keys":["requirements_payload"],"triggered_by":"user_confirmed"}

3. After packaging requirements (build_requirements_payload + write_requirements_artifact):
   emit the matching HANDOFF:: line for the destination the user chose (design or development).

NEVER say "that's outside my scope" when a user asks to proceed to development.
NEVER tell the user the pipeline order is mandatory — the user decides.
The HANDOFF:: line must appear on its own line at the very end of your response.
Do NOT add any text after it. The orchestrator strips it before showing to user.

── RULES ──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
- NEVER invent board project/team/state names — always use what tools return.
- NEVER claim a tool was called unless the tool actually returned a result in
  this turn. Do not use placeholder text such as "[Called ...]".
- NEVER include fake user messages in your answer. Do not write lines starting
  with "User:" unless quoting actual prior conversation context.
- ALWAYS confirm before any board mutation.
- Run normalize_acceptance_criteria before build_requirements_payload. Run
  detect_requirement_gaps when the user asks for requirements review or has not
  already given explicit development intent. For explicit "start development"
  requests, gaps are non-blocking risks; do not stop the handoff on gap review.
- Keep responses concise and well-formatted. Use bullet lists and headers.
- Present the completeness score and gap count prominently after gap detection.
"""

# Append the shared, agent-agnostic MCP provenance note so the agent can correctly
# identify its MCP-provided tools (namespaced `<server>__<tool>`) vs native built-ins.
INGESTION_SYS_MESSAGE = INGESTION_SYS_MESSAGE + MCP_TOOLS_PROMPT_NOTE


import logging as _logging

_req_agent_logger = _logging.getLogger(__name__)

# Per-alias orchestrator cache — avoids per-call re-construction while supporting BYOK (Pitfall 4).
_ORCHESTRATOR_CACHE: dict[tuple[str, str], object] = {}


def _build_orchestrator(model: str, litellm_provider: str, api_key: str,
                        base_url: str | None, alias: str) -> object:
    """Build (or return cached) ChatLiteLLM orchestrator. Direct-provider BYOK call.
    Cached by (alias, model) — alias is non-secret (SC#5); key never stored outside scope."""
    cache_key = (alias, model)
    if cache_key in _ORCHESTRATOR_CACHE:
        return _ORCHESTRATOR_CACHE[cache_key]
    instance = ChatLiteLLM(
        model=model,
        custom_llm_provider=litellm_provider,
        api_base=base_url,
        api_key=api_key,
        temperature=0.3,
        max_tokens=8096,
        max_retries=2,
        # Stream tokens so the copilot shows text live as it's generated (first
        # token in ~1.4s instead of a frozen wait for the whole response). Verified
        # to emit incremental chunks; LangChain still assembles the final message
        # (incl. tool calls) for the graph, so tool-calling turns are unaffected.
        streaming=True,
    )
    _ORCHESTRATOR_CACHE[cache_key] = instance
    return instance


async def agent(state: AgentState):
    from shared.services.model_resolver import resolve_model_for_run, set_resolved_model, NoModelConfiguredError, ModelNotEnabledError
    from langchain_core.messages import AIMessage
    tenant_id = state.get("tenant_id", "")
    requested = state.get("model_id")
    try:
        resolved = await resolve_model_for_run(tenant_id, requested, offering_id=state.get("offering_id"))
    except (NoModelConfiguredError, ModelNotEnabledError) as e:
        _req_agent_logger.warning("Requirements agent model resolution failed (tenant=%s): %s",
                                  tenant_id, type(e).__name__)
        return {"messages": [AIMessage(content=(
            "No usable model is configured for your organization. "
            "An administrator must add and verify a model provider in Org Settings → Model Providers."))]}
    except Exception as e:
        _req_agent_logger.exception("Requirements agent model resolution error (tenant=%s)", tenant_id)
        return {"messages": [AIMessage(content=(
            "The configured model could not be initialized. "
            "Please ask an administrator to re-verify the model provider in Org Settings → Model Providers."))]}
    set_resolved_model(resolved)
    try:
        orch = _build_orchestrator(resolved.model, resolved.litellm_provider,
                                   resolved.api_key, resolved.base_url, resolved.alias)
        # Bind base tools + any per-run MCP tools here (not in the cached builder) so
        # run-specific MCP tools never leak across runs via the shared orchestrator cache.
        orch = orch.bind_tools(tools + get_skill_tools("requirements") + get_mcp_tools())
        response = await orch.ainvoke(
            state["messages"],
            config={"metadata": {"user_api_key_alias": resolved.alias}},
        )
        # Carry the resolved model through state so the `tools` node can re-establish
        # it (its fresh asyncio.run() context loses the contextvar set above).
        return {"messages": [response], "resolved_model": resolved}
    except Exception as e:
        _req_agent_logger.exception("Requirements agent error (tenant=%s alias=%s)",
                                    tenant_id, resolved.alias)
        return {"messages": [AIMessage(content=(
            "The agent hit an error while generating a response. Please try again, "
            "or contact an administrator if the problem persists."))]}

def action(state: AgentState):

    """Synchronous action function that executes tools"""

    print("Executing tool actions...")

    last_message = state['messages'][-1]

    async def execute_tools():

        # Re-establish the BYOK-resolved model in THIS context. The sync `action`
        # node runs tools via asyncio.run(), a fresh event loop + context where the
        # `agent` node's set_resolved_model() did NOT propagate. Tools that call
        # _openai_generate (gap analysis, AC normalisation, story generation, BRD…)
        # read it via get_resolved_model(); without this they fail "No model resolved".
        # Read the already-resolved model from state (re-resolving here would hit the
        # main-loop-bound DB pool from this fresh asyncio.run() loop and fail).
        try:
            from shared.services.model_resolver import set_resolved_model, get_resolved_model
            if get_resolved_model() is None:
                _rm = state.get("resolved_model")
                if _rm is not None:
                    set_resolved_model(_rm)
        except Exception as _exc:  # noqa: BLE001 — fail-soft; tools degrade with a clear message
            _req_agent_logger.warning(
                "action: model re-establish from state failed: %s", type(_exc).__name__
            )

        results = []

        for tc in last_message.tool_calls:

            broadcast_log(manager, f"Invoking tool {tc['name']} with args {tc['args']}", level="LOGS")

            # Graceful when the model calls a tool the prompt mentions but that is
            # not bound (e.g. an unwired FLOW-G analysis tool): return an error
            # ToolMessage instead of raising StopIteration and crashing the node.
            tool_to_call = next(
                (t for t in tools + get_skill_tools("requirements") + get_mcp_tools()
                 if t.name == tc['name']), None
            )
            if tool_to_call is None:
                obs = f"Error: tool '{tc['name']}' is not available in this agent."
            else:
                # Isolate each tool: a single raising tool must NOT abort the whole batch
                # (the outer handler error-poisons every sibling tool_call). Return a
                # concise per-tool error and log the real traceback server-side.
                try:
                    obs = await tool_to_call.ainvoke(tc['args'])
                except Exception as _tool_exc:  # noqa: BLE001 — surface, don't crash the node
                    import traceback as _tb  # noqa: PLC0415
                    _req_agent_logger.error(
                        "requirements tool %s failed: %s: %s\n%s",
                        tc['name'], type(_tool_exc).__name__, _tool_exc, _tb.format_exc(),
                    )
                    _short = str(_tool_exc)
                    if len(_short) > 400:
                        _short = _short[:400] + "… (truncated — see server logs)"
                    obs = f"Tool '{tc['name']}' failed: {type(_tool_exc).__name__}: {_short}"

            # Cap before the result re-enters the LLM context (see _cap_tool_output).
            obs = _cap_tool_output(obs)
            results.append(ToolMessage(content=obs, tool_call_id=tc['id'], name=tc['name']))

            print(f"Tool {tc['name']} execution completed")

        return results

    # FIXED: Proper async execution without nested event loops

    try:

        # Check if we're already in an async context

        try:

            loop = asyncio.get_running_loop()

            # We're in an async context, so we need to handle this differently

            import nest_asyncio

            nest_asyncio.apply()

            tool_invoke_res = asyncio.run(execute_tools())

        except RuntimeError:

            # No running loop, safe to use asyncio.run

            tool_invoke_res = asyncio.run(execute_tools())

    except Exception as e:

        print(f"Error in action function: {e}")

        # Emit one error ToolMessage per pending tool_call so EVERY tool_use has a
        # matching tool_result — otherwise Anthropic rejects the next turn with
        # "unexpected tool_use_id ... has no corresponding tool_use".
        return {"messages": [
            ToolMessage(content=f"Error executing tools: {str(e)}", tool_call_id=tc["id"], name=tc["name"])
            for tc in last_message.tool_calls
        ]}

    print("All tool actions completed")

    return {"messages": tool_invoke_res}
 
 
 

def should_continue(state: AgentState) -> str:
    return "action" if state["messages"][-1].tool_calls else False

# Create workflow with synchronous functions
workflow = StateGraph(AgentState)
workflow.add_node("agent", agent)
workflow.add_node("tools", action)
workflow.add_edge("tools", "agent")
workflow.set_entry_point("agent")
workflow.add_conditional_edges("agent", should_continue, {"action": "tools", False: END})

app = workflow.compile(checkpointer=_build_checkpointer("requirements"))

# Async wrapper for the entire workflow
async def run_workflow_async(messages, config=None):

    """Async wrapper that runs the workflow using astream"""

    if config is None:

        config = {"configurable": {"thread_id": "default"}}

    print("Starting async workflow execution...")

    result = None

    async for event in app.astream({"messages": messages}, config):

        print(f"Workflow event: {list(event.keys())}")

        result = event

    print("Async workflow execution completed")

    return result

# Usage examples:

# For synchronous usage (backward compatibility):
def run_workflow_sync(messages, config=None):

    """Synchronous workflow execution using async API"""

    if config is None:

        config = {"configurable": {"thread_id": "default"}}

    print("Starting workflow execution...")

    # Use asyncio.run to run the async workflow

    result = asyncio.run(run_workflow_async(messages, config))

    print("Workflow execution completed")


    return result

# For async usage with real-time streaming:
async def main_async():
    """Example async usage"""
    messages = [HumanMessage(content="Generate a BRD from the uploaded files")]
    await run_workflow_async(messages)

# For sync usage:
def main_sync():
    """Example sync usage"""
    messages = [HumanMessage(content="Generate a BRD from the uploaded files")]
    return run_workflow_sync(messages)

# If you want to run async:
# asyncio.run(main_async())

# If you want to run sync:
# result = main_sync()
