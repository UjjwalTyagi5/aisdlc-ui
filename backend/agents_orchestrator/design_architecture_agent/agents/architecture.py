"""Design & Architecture Agent — LangGraph workflow (Azure OpenAI).

Tools:
- read_document       — extract text from uploaded local files
- generate_architecture — produce HLD/LLD/C4/API/DB/ADR from text or context
- generate_architecture_from_context — architecture from conversation context only
- update_response     — refine/update previously generated content
- save_architecture   — persist content to a .docx file
- markdowntodoc       — convert any markdown to .docx
"""
from __future__ import annotations

import asyncio
import base64
import json
import os
import pathlib
import re
import traceback
from typing import Annotated, Any, List, Optional, TypedDict

import aiofiles
import aiohttp
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import requests
from dotenv import load_dotenv
from langchain_core.messages import AIMessage, BaseMessage, HumanMessage, ToolMessage
from langchain_core.tools import tool
from config.checkpoint import build_checkpointer as _build_checkpointer
from langgraph.graph import END, StateGraph
from langgraph.graph.message import add_messages
from shared.tools.mcp_runtime import get_mcp_tools, MCP_TOOLS_PROMPT_NOTE
from shared.services.skill_runtime import get_skill_tools

from config import sdlcSettings
from config.connection_manager import manager
import logging

from config.env import (
    AGENTIC_BASE_URL,
    ANTHROPIC_MODEL,
)

logger = logging.getLogger(__name__)
from config.ws_helper import broadcast_log, get_session_id, get_user_id, get_provider_kind
from config.connectors.context import get_connector
from config.connectors.base import ConnectorNotAvailableError
from agents_orchestrator.design_architecture_agent.config import shared
from agents_orchestrator.design_architecture_agent.prompts.architecture_generation import ARCH_GEN_PROMPT
from shared.tools.spectral_tool import run_spectral_lint
from agents_orchestrator.design_architecture_agent.tools.schema_validation_tool import validate_database_schema
from agents_orchestrator.design_architecture_agent.tools.existing_system_tool import analyze_existing_system
from agents_orchestrator.design_architecture_agent.tools.figma_tools import (
    export_figma_frames,
    list_figma_frames,
    read_figma_design,
)

load_dotenv()

esett = sdlcSettings()

# Derive files directory from this file's location so it always matches
# the static server mount in process_api.py regardless of AGENTIC_APP_PATH in .env
# architecture.py → agents/ → design_architecture_agent/ → agents_orchestrator/ → agentic_app/
_FILES_DIR = str(pathlib.Path(__file__).resolve().parents[3] / "files")


async def _llm_generate_async(prompt: str, system: str = "") -> str:
    """Stream generation tokens directly to WebSocket as they arrive.

    Each token is broadcast immediately so the user sees content building in
    real time instead of waiting for the full response. Extended thinking is
    intentionally disabled — it adds 100-150 s of latency for no quality gain
    on well-specified architecture tasks.
    """
    from shared.services.model_resolver import get_resolved_model
    resolved = get_resolved_model()
    if resolved is None:
        raise RuntimeError(
            "No model resolved for this run — an administrator must configure a model provider.")

    session_id = get_session_id()
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    full_text = []
    try:
        # Deferred: importing litellm costs ~7s. sys.modules makes repeat calls free.
        import litellm
        from shared.services.model_resolver import temperature_kwargs  # noqa: PLC0415

        response = await litellm.acompletion(
            model=resolved.model,
            custom_llm_provider=resolved.litellm_provider,
            api_key=resolved.api_key,
            api_base=resolved.base_url,
            messages=messages,
            max_tokens=8192,
            # Omitted entirely for gpt-5-family models — see temperature_kwargs.
            **temperature_kwargs(resolved.model, 0.25),
            stream=True,
            timeout=120,
        )
        async for chunk in response:
            delta = chunk.choices[0].delta.content or ""
            if delta:
                full_text.append(delta)
                if session_id:
                    await manager.broadcast({
                        "type": "stream_chunk",
                        "content": delta,
                        "session_id": session_id,
                    })
    except Exception as exc:
        # If we already streamed a substantial document, a trailing/network error on the
        # closing chunk must NOT discard the whole generation (that surfaces as
        # "Tool ... failed: <entire doc>" and loses the .docx). Return what we have so it
        # can still be saved/exported; only propagate when nothing usable was produced.
        streamed = "".join(full_text)
        if len(streamed.strip()) > 200:
            logger.warning("_llm_generate_async: returning partial stream after error: %s", exc)
            broadcast_log(manager, "Generation finished with a trailing warning (content preserved).", level="INFO")
            return streamed
        broadcast_log(manager, f"Generation error: {exc}", level="ERROR")
        raise
    return "".join(full_text)


# ── File text extraction — delegated to shared utility ─────────────────────────
from shared.tools.document_tools import extract_file_text as _extract_text_from_path


# ── Markdown → DOCX helper ─────────────────────────────────────────────────────

def _add_hyperlink(paragraph, text, url):
    from docx.oxml.ns import qn
    from docx.oxml.shared import OxmlElement
    part = paragraph.part
    r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    rStyle = OxmlElement("w:rStyle")
    rStyle.set(qn("w:val"), "Hyperlink")
    rPr.append(rStyle)
    run.append(rPr)
    t = OxmlElement("w:t")
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _render_mermaid_to_png(code: str) -> bytes | None:
    """Render Mermaid diagram code to PNG bytes.

    Tries mermaid.ink first; falls back to kroki.io on any failure.
    Returns None if both services fail so the caller can degrade gracefully.
    """
    code = code.strip()
    # Primary: mermaid.ink
    try:
        encoded = base64.urlsafe_b64encode(code.encode()).decode()
        resp = requests.get(f"https://mermaid.ink/img/{encoded}?type=png", timeout=15)
        if resp.status_code == 200 and resp.content:
            return resp.content
    except Exception:
        pass
    # Fallback: kroki.io
    try:
        resp = requests.post(
            "https://kroki.io/mermaid/png",
            data=code.encode(),
            headers={"Content-Type": "text/plain"},
            timeout=20,
        )
        if resp.status_code == 200 and resp.content:
            return resp.content
    except Exception:
        pass
    return None


def _fetch_image_bytes_sync(url: str) -> bytes | None:
    """Fetch PNG/image bytes from a URL synchronously (run in executor)."""
    try:
        resp = requests.get(url, timeout=20)
        if resp.status_code == 200 and resp.content:
            return resp.content
    except Exception:
        pass
    return None


async def _fetch_image_bytes(url: str) -> bytes | None:
    """Async wrapper — fetches image from URL in a thread pool."""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _fetch_image_bytes_sync, url)


async def _markdown_to_docx(markdown_string: str, docx_path: str) -> str:
    """Thin wrapper — injects design-agent callbacks into the shared converter."""
    from shared.tools.docx_tools import markdown_to_docx as _convert
    result = await _convert(
        markdown_string,
        docx_path,
        fetch_image=_fetch_image_bytes,
        render_mermaid=_render_mermaid_to_png,
    )
    # Update shared state so the API layer can broadcast file_generated
    try:
        session_id = get_session_id()
        user_id = get_user_id()
        file_url = f"{AGENTIC_BASE_URL}/generated/{user_id}/orchestrator/{session_id}/output/{os.path.basename(docx_path)}"
        shared.output_file = os.path.basename(docx_path)
        shared.output_file_url = file_url
    except Exception:
        pass
    return result


# ── Tools ──────────────────────────────────────────────────────────────────────

@tool
def read_document(file_path: str) -> str:
    """Read a local uploaded file and return its text content.

    Supports .txt, .md, .csv, .xlsx, .xls, .docx, .pdf.
    Call this first whenever the user sends a file path.

    Args:
        file_path: Absolute path to the uploaded file.
    """
    broadcast_log(manager, f"Reading document: {file_path}", level="INFO")
    text = _extract_text_from_path(file_path)
    preview = text[:200].replace("\n", " ")
    broadcast_log(manager, f"Extracted {len(text)} chars from {os.path.basename(file_path)}: {preview}…", level="INFO")
    return text


@tool
async def generate_architecture(document_text: str, custom_prompt: str = "") -> str:
    """Generate a comprehensive architecture document (HLD/LLD/C4/API/DB schema/ADR)
    from document text or requirements.

    Args:
        document_text: Full text extracted from uploaded documents (use read_document first).
        custom_prompt: Additional focus area or instructions from the user.
    """
    broadcast_log(manager, "Generating architecture from document...", level="INFO")
    prompt_text = ARCH_GEN_PROMPT.replace("{custom_prompt}",custom_prompt or "")
    full_prompt = f"{prompt_text}\n\n--- DOCUMENT CONTENT ---\n{document_text}"
    result = await _llm_generate_async(full_prompt)
    mermaid_match = re.search(r"```mermaid(.*?)```", result, re.DOTALL)
    if mermaid_match:
        shared.mermaid = mermaid_match.group(1).strip()
    broadcast_log(manager, "Architecture generation complete.", level="INFO")
    return result


@tool
async def generate_architecture_from_context(context: str, user_requirements: str = "") -> str:
    """Generate architecture based purely on conversation context (BRD, PDD, stories, etc.)
    without requiring any uploaded file.

    Args:
        context: Complete context from conversation history (BRD, PDD, requirements, etc.).
        user_requirements: Specific additional instructions from the user.
    """
    broadcast_log(manager, "Generating architecture from context...", level="INFO")
    if not context.strip():
        return "Error: Context cannot be empty."

    system_msg = (
        "You are a Senior AI Solutions Architect. "
        "You MUST produce a complete enterprise-grade architecture document following the template exactly. "
        "NEVER skip a required section. ALWAYS include valid Mermaid code blocks for diagrams. "
        "DB schema MUST be full CREATE TABLE SQL. API contracts MUST include request/response JSON."
    )
    prompt = f"""{ARCH_GEN_PROMPT.replace("{custom_prompt}",user_requirements or "")}

--- CONVERSATION CONTEXT (use this as the source of requirements) ---
{context}

IMPORTANT: Replace every placeholder in the template above with content derived from the
conversation context. All three C4 Mermaid diagrams are REQUIRED. DB schema SQL is REQUIRED.
Do not leave any table cells empty in required sections."""

    result = await _llm_generate_async(prompt, system_msg)
    # Stash the generated markdown so save_architecture can persist it even when a
    # weaker model fails to echo the (large) document back into the save tool's
    # `content` argument — the common cause of an empty/failed .docx.
    shared.last_architecture = result
    mermaid_match = re.search(r"```mermaid(.*?)```", result, re.DOTALL)
    if mermaid_match:
        shared.mermaid = mermaid_match.group(1).strip()
    broadcast_log(manager, "Context-based architecture generation complete.", level="INFO")
    return result


@tool
async def update_response(query: str, content: str, file_paths: Optional[List[str]] = None) -> str:
    """Update/refine previously generated content based on a query.

    Args:
        query: The update request or instructions from the user.
        content: The content to be updated.
        file_paths: Optional list of local file paths for additional context.
    """
    broadcast_log(manager, "Updating response...", level="INFO")
    extra_context = ""
    if file_paths:
        texts = [_extract_text_from_path(p) for p in file_paths if p]
        extra_context = "\n\n--- ADDITIONAL FILE CONTEXT ---\n" + "\n\n".join(texts)

    prompt = f"""Update the content below based on the query. Preserve formatting and style.

QUERY: {query}

CONTENT TO UPDATE:
{content}{extra_context}

Return only the updated content."""
    result = await _llm_generate_async(prompt)
    broadcast_log(manager, "Response updated.", level="INFO")
    return result


@tool
async def markdowntodoc(content: str, output_path: str) -> str:
    """Convert a Markdown string to a .docx file and save it.

    Args:
        content: Markdown string to convert.
        output_path: Relative output path, e.g. 'architecture.docx'.
    """
    session_id = get_session_id()
    user_id = get_user_id()
    out_dir = os.path.join(_FILES_DIR, str(user_id), "orchestrator", str(session_id), "output")
    os.makedirs(out_dir, exist_ok=True)
    full_path = os.path.join(out_dir, os.path.basename(output_path))
    return await _markdown_to_docx(content, full_path)


@tool
async def save_architecture(filename: str = "", content: str = "") -> str:
    """Save the generated architecture document to a .docx file.

    Call this whenever the user says 'save', 'export', or 'download'.
    Always pass the full architecture markdown as 'content'.

    Args:
        filename: Output filename, e.g. 'project_architecture.docx'.
                  Defaults to 'architecture.docx' if omitted.
        content:  The full architecture markdown text to save. Required.
    """
    if not content or not content.strip():
        # Fallback to the last architecture generated this session — a weaker model
        # often cannot copy the full markdown into this argument.
        content = getattr(shared, "last_architecture", "") or ""
    if not content or not content.strip():
        return (
            "ERROR: No architecture content available. "
            "Generate the architecture first with generate_architecture_from_context, "
            "then call save_architecture again."
        )
    if not filename or not filename.strip():
        filename = "architecture.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"
    broadcast_log(manager, f"Saving architecture document: {filename}", level="INFO")
    session_id = get_session_id()
    user_id = get_user_id()
    out_dir = os.path.join(_FILES_DIR, str(user_id), "orchestrator", str(session_id), "output")
    full_path = os.path.join(out_dir, filename)
    result = await _markdown_to_docx(content, full_path)
    broadcast_log(manager, f"Architecture document saved: {filename}", level="INFO")
    # Broadcast file_generated directly from the tool so the saved doc reliably
    # surfaces on the main screen (the post-stream WS-handler broadcast depends on
    # shared.output_file surviving across the tools-node task context, which it may
    # not). Mirrors the requirements agent's behaviour.
    _dl_url = ""
    try:
        file_size = os.path.getsize(full_path) if os.path.exists(full_path) else 0
        _dl_url = getattr(shared, "output_file_url", "") or f"{AGENTIC_BASE_URL}/generated/{user_id}/orchestrator/{session_id}/output/{os.path.basename(full_path)}"
        await manager.broadcast({
            "type": "file_generated",
            "session_id": session_id,
            "filename": os.path.basename(full_path),
            "url": _dl_url,
            "file_size": file_size,
            "agent_name": "Design Agent",
        })
    except Exception as _exc:  # noqa: BLE001 — surfacing is best-effort
        logger.warning("save_architecture: file_generated broadcast failed: %s", _exc)
    # Persist the saved .docx as a project Artifact so it shows in the artifacts panel.
    try:
        from shared.services.chat_artifacts import register_generated_file  # noqa: PLC0415
        await register_generated_file(os.path.basename(full_path), full_path, _dl_url, stage="design")
    except Exception:  # noqa: BLE001 — best-effort
        logger.debug("save_architecture register_generated_file failed", exc_info=True)
    # Return the real download URL so the agent presents a working link (not a fabricated one).
    if _dl_url:
        return f"Saved '{filename}'. Download it here: {_dl_url}"
    return result


@tool
async def render_diagram_via_kroki(diagram_type: str, source: str) -> str:
    """Render any diagram as an SVG image URL via kroki.io.

    Use this for diagram types that Mermaid does not support natively, such as
    PlantUML, GraphViz (dot), BPMN, or Excalidraw. The returned URL can be
    embedded in markdown as an image: ![diagram](url)

    Args:
        diagram_type: One of: plantuml, graphviz, bpmn, excalidraw, c4plantuml,
                      ditaa, erd, pikchr, vega, vegalite, wavedrom.
        source: The raw diagram source code in the appropriate syntax.

    Returns a public kroki.io URL pointing to the rendered SVG.
    """
    import base64 as _b64
    import zlib as _zlib

    diagram_type = diagram_type.strip().lower()
    compressed = _b64.urlsafe_b64encode(
        _zlib.compress(source.encode("utf-8"), 9)
    ).decode("utf-8")
    # Use PNG so the image can be embedded directly in docx exports
    png_url = f"https://kroki.io/{diagram_type}/png/{compressed}"
    # Also provide SVG for browser rendering (higher quality)
    svg_url = f"https://kroki.io/{diagram_type}/svg/{compressed}"
    broadcast_log(manager, f"Kroki URL generated for {diagram_type}", level="INFO")
    # Return a markdown image using the PNG URL — _markdown_to_docx will fetch
    # and embed this PNG automatically when saving to .docx
    return (
        f"![{diagram_type} diagram]({png_url})\n\n"
        f"SVG (browser): {svg_url}"
    )


@tool
async def update_ado_epic_design_complete(
    project: str,
    epic_id: int,
    design_url: str,
    target_state: str = "Design Complete",
) -> str:
    """Mark the parent ADO epic as Design Complete and attach the design document URL.

    Call this after save_architecture succeeds AND the user confirms design is complete.
    Always confirm with the user before calling this tool.

    Args:
        project: ADO project name.
        epic_id: Work item ID of the parent epic.
        design_url: URL of the generated design document (from the save_architecture result).
        target_state: State to move the epic to. Default 'Design Complete'.
                      Use list_ado_states if unsure of the exact state name.
    """
    # CONSEQUENTIAL (§1.5): this moves a real epic and comments on it in front of the
    # team. The docstring above used to be the only control — "Always confirm with the
    # user before calling this tool" is a request to the model, not an enforced rule,
    # and the tool node executes whatever the model emits. Same shared rule the
    # Requirements board writes go through.
    # BOTH halves of the rule: the Architect's role AND their explicit yes on this
    # turn. The role check alone would pass every time an Architect drove the chat,
    # which is exactly who drives it — so the epic would move with nobody asked.
    from shared.authz.consequential import authorize_consequential  # noqa: PLC0415 — import cycle

    ok, why = await authorize_consequential(
        "design",
        action=f"Moving epic #{epic_id} to '{target_state}' and commenting on it",
        ask=(
            f"ask the user \"Shall I move epic #{epic_id} to '{target_state}' and post "
            f"the design document link on it?\""
        ),
    )
    if not ok:
        return why

    try:
        connector = get_connector()
    except RuntimeError:
        return (
            "No project-management board is connected for this run, so the epic could "
            "not be updated. An administrator can connect one on the Integrations page."
        )

    results = []
    try:
        await connector.write_adapter("move_item_state", project=project, item_id=epic_id, new_state=target_state)
        results.append(f"Epic #{epic_id} moved to '{target_state}'.")
    except Exception as exc:  # noqa: BLE001
        # The TYPE only. A connector error carries the instance URL and the full API
        # path, and this string goes straight into the model's context and into the
        # saved transcript. ConnectorAccessDenied is named separately because its fix
        # is an access level, not a retry.
        name = type(exc).__name__
        results.append(
            f"Could not move the epic state ({name})."
            + (" This project's design stage does not have write access to its board;"
               " an administrator can change that on the Integrations page."
               if name == "ConnectorAccessDenied"
               else f" Check that '{target_state}' is a real state on this board.")
        )

    try:
        comment = f"Design complete. Architecture document: {design_url}"
        await connector.write_adapter("add_comment", project=project, item_id=epic_id, comment=comment)
        results.append(f"Comment with design document link added to #{epic_id}.")
    except Exception as exc:  # noqa: BLE001
        results.append(f"Could not add the comment ({type(exc).__name__}).")

    return "\n".join(results)


@tool
async def write_design_artifact(
    run_id: str,
    hld: str = "",
    lld: str = "",
    api_contracts: str = "",
    database_schema: str = "",
    c4_diagram_url: str = "",
    security_checklist: str = "",
) -> str:
    """Persist the design artifact to the database and notify the orchestrator.

    Call this after the user confirms the design is complete to advance the
    pipeline to the development stage.

    Args:
        run_id: Pipeline run identifier (provided in the system context).
        hld: High-level design section text.
        lld: Low-level design section text.
        api_contracts: API contract section text.
        database_schema: Database schema section text.
        c4_diagram_url: Optional URL of the rendered C4 diagram.
        security_checklist: OWASP Top-10 security design checklist section text.
    """
    from shared.services.artifact_service import write_and_notify
    from shared.models.artifacts import DesignArtifact

    artifact = DesignArtifact(
        hld=hld or None,
        lld=lld or None,
        api_contracts=api_contracts or None,
        database_schema=database_schema or None,
        c4_diagram_url=c4_diagram_url or None,
        security_checklist=security_checklist or None,
    )
    await write_and_notify(run_id, "design", artifact.model_dump())
    return f"Design artifact persisted for run {run_id}"


# ── Downloadable-artifact helpers (docx already handled by save_architecture) ─────

def _design_output_dir() -> str:
    out = os.path.join(_FILES_DIR, str(get_user_id()), "orchestrator", str(get_session_id()), "output")
    os.makedirs(out, exist_ok=True)
    return out


async def _design_broadcast_file(filename: str, path: str) -> str:
    """Broadcast a file_generated event for a produced file and return its download URL."""
    try:
        user_id = get_user_id()
        session_id = get_session_id()
        size = os.path.getsize(path) if os.path.exists(path) else 0
        url = f"{AGENTIC_BASE_URL}/generated/{user_id}/orchestrator/{session_id}/output/{filename}"
        await manager.broadcast({
            "type": "file_generated",
            "session_id": session_id,
            "filename": filename,
            "url": url,
            "file_size": size,
            "agent_name": "Design Agent",
        })
        try:
            from shared.services.chat_artifacts import register_generated_file  # noqa: PLC0415
            await register_generated_file(filename, path, url, stage="design")
        except Exception:  # noqa: BLE001 — best-effort; never block generation
            logger.debug("design register_generated_file failed", exc_info=True)
        return url
    except Exception as exc:  # noqa: BLE001 — surfacing is best-effort
        logger.warning("design file_generated broadcast failed: %s", exc)
        return ""


def _design_download_bytes(url: str) -> Optional[bytes]:
    try:
        import httpx as _httpx  # noqa: PLC0415
        r = _httpx.get(url, timeout=30, follow_redirects=True)
        r.raise_for_status()
        return r.content
    except Exception:
        return None


@tool
async def generate_ppt(content: str, output_path: str) -> str:
    """Generate a PowerPoint (.pptx) deck from markdown and make it downloadable from chat —
    e.g. an architecture review deck. '#'/'##' headings start slides; '- ' lines are bullets.

    Args:
        content: Markdown-style deck content (headings = slides, bullets = points).
        output_path: Filename, e.g. 'architecture_overview.pptx'.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        return "Error: PowerPoint generation needs python-pptx — run `pip install python-pptx` on the backend."
    filename = os.path.basename(output_path) or "presentation.pptx"
    if not filename.lower().endswith(".pptx"):
        filename += ".pptx"
    dest = os.path.join(_design_output_dir(), filename)

    def _build() -> None:
        prs = Presentation()
        slides: List[tuple] = []
        cur_title: Optional[str] = None
        cur_bullets: List[str] = []
        for raw in content.splitlines():
            s = raw.strip()
            if not s:
                continue
            if s.startswith("#"):
                if cur_title is not None:
                    slides.append((cur_title, cur_bullets))
                cur_title, cur_bullets = s.lstrip("#").strip(), []
            elif s.startswith(("- ", "* ")):
                cur_bullets.append(s[2:].strip())
            else:
                cur_bullets.append(s)
        if cur_title is not None:
            slides.append((cur_title, cur_bullets))
        if not slides:
            slides = [("Presentation", [content.strip()[:400]])]
        ft, fb = slides[0]
        ts = prs.slides.add_slide(prs.slide_layouts[0])
        ts.shapes.title.text = ft
        if fb and len(ts.placeholders) > 1:
            ts.placeholders[1].text = fb[0]
        for title, bullets in slides[1:]:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets or [""]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
        prs.save(dest)

    await asyncio.get_event_loop().run_in_executor(None, _build)
    url = await _design_broadcast_file(filename, dest)
    tail = f" Download it here: {url}" if url else ""
    return f"Created PowerPoint deck '{filename}'.{tail}"


@tool
async def generate_diagram(diagram_type: str, source: str, output_path: str = "") -> str:
    """Render a diagram as a downloadable PNG file (via kroki.io) the user can save from chat.
    Use this when the user wants to DOWNLOAD an image; use render_diagram_via_kroki when you
    only need to EMBED the diagram inline in the architecture document.

    Args:
        diagram_type: mermaid, plantuml, graphviz, c4plantuml, erd, bpmn, excalidraw, ...
        source: Diagram source code in that syntax.
        output_path: Optional filename, e.g. 'component_diagram.png'.
    """
    import base64 as _b64  # noqa: PLC0415
    import zlib as _zlib  # noqa: PLC0415

    dtype = diagram_type.strip().lower()
    compressed = _b64.urlsafe_b64encode(_zlib.compress(source.encode("utf-8"), 9)).decode("utf-8")
    png_url = f"https://kroki.io/{dtype}/png/{compressed}"
    img = await asyncio.get_event_loop().run_in_executor(None, _design_download_bytes, png_url)
    if not img:
        return f"Error: could not render the {dtype} diagram — check the type/syntax."
    filename = os.path.basename(output_path) or f"diagram_{dtype}.png"
    if not filename.lower().endswith(".png"):
        filename += ".png"
    dest = os.path.join(_design_output_dir(), filename)
    with open(dest, "wb") as fh:
        fh.write(img)
    url = await _design_broadcast_file(filename, dest)
    tail = f" Download it here: {url}" if url else ""
    return f"Rendered {dtype} diagram '{filename}'.{tail}"



@tool
async def save_architecture_pdf(content: str = "", filename: str = "architecture.pdf") -> str:
    """Save the architecture document as a PDF instead of a .docx.

    Use when the user asks for a PDF. For Word use save_architecture; for slides use
    generate_ppt.

    Args:
        content: The markdown document. Omit to use the most recently generated
            architecture for this session.
        filename: Output filename, e.g. project_architecture.pdf
    """
    if not (content and content.strip()):
        # `last_architecture`, NOT `output_file`. output_file holds a FILENAME
        # (os.path.basename(docx_path)) — reading it as content produced
        # "Error: nothing to save" for a user who had a full document on screen,
        # because the fallback was looking at the wrong attribute entirely.
        content = getattr(shared, "last_architecture", "") or ""
    if not (content and content.strip()):
        return ("Error: nothing to save. Generate the architecture first "
                "(generate_architecture_from_context), or pass the content explicitly.")

    filename = os.path.basename(filename or "architecture.pdf")
    if not filename.lower().endswith(".pdf"):
        filename += ".pdf"
    user_id, session_id = get_user_id(), get_session_id()
    out_dir = os.path.join(_FILES_DIR, str(user_id), "orchestrator", str(session_id), "output")
    os.makedirs(out_dir, exist_ok=True)
    full_path = os.path.join(out_dir, filename)

    try:
        from shared.tools.pdf_render import markdown_to_pdf  # noqa: PLC0415

        markdown_to_pdf(content, full_path, title=filename.rsplit(".", 1)[0])
    except Exception as exc:  # noqa: BLE001
        return f"Error generating the PDF ({type(exc).__name__})."

    url = await _design_broadcast_file(filename, full_path)
    return (
        f"Saved '{filename}'."
        + (f" Download it here: {url}" if url else "")
        + " It is NOT yet in the project's artifacts — ask the user whether to save it "
          "there. It is awaiting a project admin's approval."
    )


@tool
async def read_project_requirements() -> str:
    """Load this project's requirements — the stories, their descriptions and their
    acceptance criteria — from the Requirements stage.

    Call this when the user asks you to design something FOR THIS PROJECT, or refers to
    "the requirements", "the stories", "the backlog" or similar. Do NOT call it when the
    user is asking you to design something they have described themselves in the chat:
    their words are the requirement then, and loading the project's backlog on top of
    them would design the wrong system.

    Returns the requirements as text, or says plainly that the project has none.

    THIS REPLACED AN AUTOMATIC INJECTION. The Design page used to push every story into
    the agent's context before the user had typed anything, behind a "From requirements
    / Freeform" toggle that defaulted to on. That fed the agent whatever the board held
    — Epics and project-setup Tasks included — and it designed a system for them. Making
    it a tool is what `_build_session_context` argued for all along: context the model
    chooses to load, not an injection it cannot decline.
    """
    from config.context_broker import build_context_for_project  # noqa: PLC0415
    from config.ws_helper import get_project_id, get_tenant_id  # noqa: PLC0415

    project_id, tenant_id = get_project_id(), get_tenant_id()
    if not project_id or not tenant_id:
        # A standalone conversation outside any project. Say so rather than returning
        # an empty string the model would read as "the project has no requirements".
        return (
            "This conversation is not attached to a project, so there are no stored "
            "requirements to read. Design from what the user has described."
        )

    context = await build_context_for_project(project_id, tenant_id, "design")
    if not context:
        return (
            "This project has no requirements recorded yet. Design from what the user "
            "has described, or ask them what the system needs to do."
        )
    return context



@tool
async def export_document(content: str = "", filename: str = "architecture.docx") -> str:
    """Export content as a Word (.docx), PDF (.pdf), Excel (.xlsx), Markdown (.md) or text file.

    ONE tool for every document format — the format comes from the FILENAME EXTENSION,
    so "as a PDF" means filename="something.pdf".

        Word         architecture.docx
        PDF          architecture.pdf
        Excel        api_contract.xlsx  (exports the MARKDOWN TABLES, one sheet per table)
        Markdown     architecture.md

    Args:
        content: The markdown to export. Omit to use the most recently generated
            architecture document for this session.
        filename: Output filename INCLUDING the extension you want.

    For diagrams/images use generate_diagram or render_diagram_via_kroki; for slides
    use generate_ppt. save_architecture / save_architecture_pdf remain for the plain
    architecture document; this tool exports ANY content in ANY of these formats.
    """
    from shared.tools.doc_export import (  # noqa: PLC0415
        export_result_message,
        normalise_filename,
        render_document,
        supported_list,
    )

    if not (content and content.strip()):
        # `last_architecture`, NOT `output_file`. output_file holds a FILENAME
        # (os.path.basename(docx_path)) — reading it as content produced
        # "Error: nothing to save" for a user who had a full document on screen,
        # because the fallback was looking at the wrong attribute entirely.
        content = getattr(shared, "last_architecture", "") or ""
    if not (content and content.strip()):
        return ("Error: nothing to export. Generate the architecture first "
                "(generate_architecture_from_context), or pass the content explicitly.")

    name = normalise_filename(filename, "architecture.docx")
    user_id, session_id = get_user_id(), get_session_id()
    out_dir = os.path.join(_FILES_DIR, str(user_id), "orchestrator", str(session_id), "output")
    os.makedirs(out_dir, exist_ok=True)
    full_path = os.path.join(out_dir, name)

    try:
        await render_document(content, full_path, title=name.rsplit(".", 1)[0])
    except ValueError:
        # render_document refuses an extension it has no renderer for. Its message is
        # ours, not a connector's, but the sweep in test_board_write_failures cannot
        # tell those apart from source text — and the filename already names the
        # offending extension, so interpolating the exception adds nothing.
        return f"Error: '{name}' has an unsupported extension. Supported: {supported_list()}"
    except Exception as exc:  # noqa: BLE001
        return f"Error generating '{name}' ({type(exc).__name__}). Supported: {supported_list()}"

    url = await _design_broadcast_file(name, full_path)
    extras = []
    if name.lower().endswith(".xlsx"):
        extras.append("Excel exports contain the document's TABLES, one sheet each.")
    return export_result_message(name, url, extras)

tools = [
    read_document,
    # The project's requirements, ON DEMAND. Replaced the Design page's automatic
    # injection, which pushed every board item into context before the user had typed.
    read_project_requirements,
    generate_architecture,
    generate_architecture_from_context,
    update_response,
    # markdowntodoc intentionally NOT bound — it's redundant with save_architecture
    # (which converts markdown→docx) and a weaker model fixates on it, calling it with
    # empty content and looping. Removing it forces the correct
    # generate_architecture_from_context → save_architecture path.
    save_architecture,
    # PDF output, and the explicit save the user is asked for before anything is
    # written to the project's shared artifact storage.
    save_architecture_pdf,
    export_document,
    generate_ppt,
    generate_diagram,
    render_diagram_via_kroki,
    update_ado_epic_design_complete,
    write_design_artifact,
    analyze_existing_system,
    run_spectral_lint,
    validate_database_schema,
    # Figma — a design that already exists is a better input than a described one.
    list_figma_frames,
    read_figma_design,
    export_figma_frames,
]


# ── System prompt ──────────────────────────────────────────────────────────────

# WHY "ACT, DON'T NARRATE" IS WORDED THE WAY IT IS. It used to end with "if the
# structured pipeline context already contains requirements / user stories ... call
# generate_architecture_from_context immediately". The presence of context WAS the
# instruction, so a user who opened the Design page and typed "hi" received a complete
# eight-section architecture document they had never asked for. The Design page threads
# the project's stories into pipeline_context by default ("From requirements" mode), so
# this fired on the first message of every project that had requirements.
#
# The anti-narration rule itself is worth keeping — a model that says "I'll generate
# that now" and then stops is a real failure. The trigger just has to be the user
# asking. See tests/test_design_greeting_does_not_generate.py.
DESIGN_SYS_MESSAGE = """\
You are the Design & Architecture Agent — powered by Claude. You transform
requirements, BRDs, PDDs, and user stories into comprehensive, enterprise-grade
architecture documents structured across 8 standard artifacts.

── ACT, DON'T NARRATE (CRITICAL) ─────────────────────────────────────────────
When the user ASKS you to generate, design, or save, you MUST emit the tool call
in the SAME response — do NOT reply with only a preamble like "I'll generate it
now" and then stop. A turn that announces an action without calling the tool is a
failure.

THE USER ASKING IS THE TRIGGER — NOT THE PRESENCE OF CONTEXT. Structured pipeline
context (requirements / user stories) is SOURCE MATERIAL for when you are asked to
design. It is not an instruction, and it is not consent. If the user greets you,
asks a question, or says something you cannot read as a request to produce a design,
REPLY TO WHAT THEY SAID. Say briefly what you can see (e.g. "I can see N stories
from Requirements") and ask what they want built.

When you ARE asked, use the context you have: pass those stories as `context`, and
never ask the user to re-supply requirements that are already in front of you.

── THE PROJECT'S REQUIREMENTS ARE A TOOL CALL, NOT SOMETHING YOU ARE HANDED ──────
You do NOT automatically receive this project's stories. Call
`read_project_requirements` to load them.

CALL IT when the user asks you to design something FOR THIS PROJECT, or refers to
"the requirements", "the stories", "the backlog", "what we captured" or similar.
Call it BEFORE designing in that case — designing from memory when stored
requirements exist produces a system nobody asked for.

DO NOT CALL IT when the user describes what they want in the chat. Their words ARE
the requirement then. Loading the project's backlog on top of a description the user
just gave you designs the wrong system — and a project's board holds Epics and
setup Tasks alongside real stories, none of which are things to design.

DO NOT CALL IT on a greeting. "hi" is not a request to design anything.


── SAVING DOCUMENTS TO THE PROJECT (AUTOMATIC, THEN APPROVED) ────────────────────
Every document you generate is recorded in the project's artifacts automatically, as
AWAITING APPROVAL. You do NOT ask whether to save it, and there is no tool to call.
- After generating a document, tell the user it is ready and give the download link.
- Then say it has been submitted for approval, and that a project admin decides whether
  it joins the project's shared record.
- Do NOT claim it has been "added to the project" or "saved to artifacts" — it is
  waiting on someone else's decision, and saying otherwise sets the wrong expectation.
- The download link works immediately either way.

── FILE FORMATS ──────────────────────────────────────────────────────────────────
export_document produces EVERY document format — the format comes from the FILENAME
EXTENSION you pass:
  Word         export_document(filename="report.docx")
  PDF          export_document(filename="report.pdf")
  Excel        export_document(filename="matrix.xlsx")  — exports the document's
               TABLES, one sheet per table (not prose)
  Markdown     export_document(filename="notes.md")
  Diagram      generate_diagram
  Slides       generate_ppt
Never substitute a format silently. If the user asks for PDF, pass a .pdf filename.
If they ask for a format not listed here, say which ones are available.

── SCOPE BOUNDARY (CRITICAL) ─────────────────────────────────────────────────
You are ONLY responsible for design and architecture artifacts.
- If asked about code, implementation, or development: say "That's handled by the
  Development Agent." then call write_design_artifact to persist your output.
- If asked about test plans or QA: say "That's handled by the Testing Agent."
  then call write_design_artifact to persist your output if design is complete.
- If asked for a REQUIREMENTS artifact — a BRD, PDD, MoM, risk register, user
  stories, acceptance criteria, or a requirements document of any kind — say
  "That's handled by the Requirements Agent, on the Requirements page." and STOP.
  You have no tool that produces any of them.

NEVER SUBSTITUTE A DIFFERENT ARTIFACT FOR THE ONE ASKED FOR. If you cannot produce
what the user named, say which agent can and stop. Producing an architecture
document because a BRD was requested wastes minutes of the user's time, spends
their tokens, and answers a question nobody asked — the user has to read the whole
thing to discover it is not what they wanted. "I can't do that here, the
Requirements Agent can" is a better answer than a document.

── ACCURACY & GROUNDING (READ FIRST — CRITICAL FOR DEMO) ────────────────────
These rules prevent hallucination. Every section of the architecture document
must be derived from the requirements provided — not invented.

1. REQUIREMENTS FIRST. Before generating architecture, you MUST have one of:
   - Uploaded file(s) the user provided, OR
   - A requirements payload injected from the Requirements Agent, OR
   - The user's direct description of what to build.
   If none of these exist, ask: "Could you share the requirements or user
   stories before I start? This ensures the design matches your actual system."

2. NO SILENT INVENTION. If a requirement, entity, endpoint, or business rule
   is not in the source material, do NOT silently include it. Either:
   - Ask the user first, or
   - Include it tagged as: > ⚠️ **[ASSUMPTION]** — not specified in requirements.
     Please confirm or correct this before development begins.

3. API ENDPOINTS must map to user stories. For each endpoint in the API
   CONTRACT section, include a one-line annotation:
   "Implements: [User Story / Feature name from requirements]"
   Never define endpoints for features not mentioned in requirements.

4. DATABASE TABLES must map to domain entities from requirements. If you name
   a table or column that wasn't in the source, tag it [ASSUMPTION].

5. IF CONTEXT IS THIN → ASK, DON'T GUESS.
   If requirements are vague (e.g. "build a leave management app" with no
   stories), ask 2–3 targeted questions before calling generate_architecture:
   - "What are the main user roles?"
   - "What are the key actions each role performs?"
   - "Are there any existing systems this must integrate with?"

── WORKFLOW ──────────────────────────────────────────────────────────────────
0. EXISTING-SYSTEM ANALYSIS (do this FIRST when a repo path is available):
   call analyze_existing_system(repo_path) to detect the current frameworks,
   language mix, and API/DB files. Design to EXTEND the existing system — reuse
   its stack and patterns instead of proposing a greenfield rewrite. If no repo
   path is available, skip this step.
0b. EXISTING DESIGN (do this when the user mentions Figma or gives a figma.com link):
   call list_figma_frames(file_url) to get the screen inventory, then
   read_figma_design(file_url, node_ids) on the screens that matter. Design to the
   screens that EXIST — name the real components and reuse the real copy instead of
   inventing a UI the designers did not draw. Use export_figma_frames(node_ids) only
   when the user wants images embedded in the document; its URLs expire in ~30 days,
   so never present them as permanent. If Figma is not connected, the tools say so —
   carry on from the written requirements rather than stopping.
1. If the user provides file paths → call read_document for EACH file first.
2. Pass extracted text to generate_architecture.
3. If no files → use generate_architecture_from_context with conversation context.
4. VALIDATION LOOPS (run after generation, before asking to save):
   • API CONTRACT — save the OpenAPI YAML to a file, then call
     run_spectral_lint(spec_path). If it returns findings with severity "error",
     fix the spec and re-lint until clean (or "unavailable"). Advisory only —
     proceed if the tool is unavailable.
   • DATABASE SCHEMA — call validate_database_schema(ddl) on the CREATE TABLE
     block. For each issue (e.g. missing_primary_key), correct the DDL and
     re-validate. Advisory only — proceed if unavailable.
   These tools degrade gracefully — never block the design if one is unavailable.
5. After generation + validation, ask if the user wants to save as .docx (save_architecture).
6. On "save", "export", "download" → call save_architecture with ONLY the filename.
   NEVER pass 'content' — the system caches the last generated architecture automatically.
   Example: save_architecture(filename="project_architecture.docx")
7. On update/refinement requests → call update_response.
   When updating, preserve ALL 8 section headers — never drop a section when
   refining only one part of the document.

── OUTPUT FORMAT (CRITICAL — frontend parses these exact headers) ─────────────
Structure EVERY architecture response with ALL 8 sections using these EXACT
headers (uppercase, no trailing punctuation):

## HIGH-LEVEL DESIGN
## LOW-LEVEL DESIGN
## C4 ARCHITECTURE DIAGRAMS
## API CONTRACT
## DATABASE SCHEMA
## ARCHITECTURE DECISION RECORDS
## TECHNOLOGY STACK
## SECURITY DESIGN CHECKLIST

The frontend splits on these exact headers to render tabbed artifact panels.
If a section is not yet applicable, include the header with a one-line note.

── ARTIFACT TEMPLATES ────────────────────────────────────────────────────────

### HIGH-LEVEL DESIGN
1. System Overview (2–3 paragraphs)
2. Component Architecture — MANDATORY Mermaid flowchart:
   ```mermaid
   graph TD
       A[Browser / Mobile] --> B[React SPA]
       B --> C[API Gateway]
       C --> D[Service Layer]
       D --> E[(Database)]
   ```
3. Data Flow — Mermaid sequence diagram for the primary use case
4. Integration Points — table: System | Protocol | Auth | Purpose
5. NFR Summary — table: Category | Requirement | Target

### LOW-LEVEL DESIGN
1. Component Specifications — per component: responsibilities, interfaces, dependencies
2. Component / Class Diagram — MANDATORY Mermaid `classDiagram` (or detailed
   `flowchart` if classes don't fit the domain) showing concrete classes/modules,
   their key methods/fields, and relationships. NEVER omit this diagram.
3. Sequence Diagrams — MANDATORY Mermaid sequence diagram for each key flow
   ```mermaid
   sequenceDiagram
       participant U as User
       participant API as REST API
       participant DB as Database
       U->>API: GET /resource
       API->>DB: SELECT query
       DB-->>API: Result rows
       API-->>U: 200 JSON response
   ```
4. Class / Module Structure — class names, methods, data types per service
5. Error Handling Strategy — error codes, retry policy, fallback behaviour

### C4 ARCHITECTURE DIAGRAMS
Use Mermaid C4 syntax for all three levels:

Level 1 — Context:
```mermaid
C4Context
    Person(user, "User", "Description")
    System(sys, "System", "Description")
    Rel(user, sys, "Uses", "HTTPS")
```

Level 2 — Container:
```mermaid
C4Container
    Container(web, "Web App", "React", "SPA")
    Container(api, "API", "FastAPI", "REST")
    ContainerDb(db, "Database", "PostgreSQL", "Stores data")
    Rel(web, api, "Calls", "HTTPS/JSON")
    Rel(api, db, "Reads/Writes", "SQL")
```

Level 3 — Component:
```mermaid
C4Component
    Component(ctrl, "Controller", "FastAPI", "Handles requests")
    Component(svc, "Service", "Python", "Business logic")
    Component(repo, "Repository", "SQLAlchemy", "Data access")
    Rel(ctrl, svc, "Calls")
    Rel(svc, repo, "Uses")
```

### API CONTRACT
Produce full OpenAPI 3.0 YAML in a fenced yaml block:
```yaml
openapi: 3.0.3
info:
  title: <Service> API
  version: 1.0.0
paths:
  /api/resource:
    get:
      summary: ...
      parameters: [...]
      responses:
        '200':
          description: ...
          content:
            application/json:
              schema:
                $ref: '#/components/schemas/ResourceResponse'
components:
  schemas:
    ResourceResponse:
      type: object
      properties:
        id: { type: integer }
        name: { type: string }
```

### DATABASE SCHEMA
Include ER diagram and DDL:
```mermaid
erDiagram
    TABLE_A {
        int id PK
        string name
    }
    TABLE_B {
        int id PK
        int table_a_id FK
    }
    TABLE_A ||--o{ TABLE_B : "has"
```
```sql
CREATE TABLE table_a (
    id SERIAL PRIMARY KEY,
    name VARCHAR(255) NOT NULL
);
CREATE TABLE table_b (
    id SERIAL PRIMARY KEY,
    table_a_id INTEGER NOT NULL REFERENCES table_a(id)
);
CREATE INDEX idx_table_b_a ON table_b(table_a_id);
```

### ARCHITECTURE DECISION RECORDS
For each significant decision:

**ADR-001: [Decision Title]**
- Status: Accepted
- Context: [Why a decision was needed]
- Decision: [What was decided]
- Rationale: [Why this option over alternatives]
- Consequences: [Trade-offs, risks, impacts]
- Alternatives Rejected: [Other options considered]

### TECHNOLOGY STACK
| Layer | Technology | Version | Justification |
|-------|-----------|---------|---------------|
| Frontend | ... | ... | ... |
| Backend | ... | ... | ... |
| Database | ... | ... | ... |
| Auth | ... | ... | ... |
| Hosting | ... | ... | ... |
| CI/CD | ... | ... | ... |

### SECURITY DESIGN CHECKLIST
Review the designed APIs and data model against the OWASP Top 10 (2021). For EACH
category, state the risk in the context of THIS design and the mitigating control:

| OWASP Category | Applicable? | Risk in this design | Mitigation control |
|----------------|-------------|---------------------|--------------------|
| A01 Broken Access Control | ... | ... | ... |
| A02 Cryptographic Failures | ... | ... | ... |
| A03 Injection | ... | ... | ... |
| A04 Insecure Design | ... | ... | ... |
| A05 Security Misconfiguration | ... | ... | ... |
| A06 Vulnerable & Outdated Components | ... | ... | ... |
| A07 Identification & Authentication Failures | ... | ... | ... |
| A08 Software & Data Integrity Failures | ... | ... | ... |
| A09 Security Logging & Monitoring Failures | ... | ... | ... |
| A10 Server-Side Request Forgery (SSRF) | ... | ... | ... |
Tag any control you ASSUME (not stated in requirements) with ⚠️ [ASSUMPTION].

── DIAGRAM RENDERING ─────────────────────────────────────────────────────────
Primary: Use Mermaid (```mermaid blocks) for all standard diagrams — flowcharts,
sequence, ER, class, C4Context/C4Container/C4Component.

Fallback (render_diagram_via_kroki): Use when the user asks for diagram types
Mermaid cannot render natively:
- PlantUML UML diagrams         → diagram_type="plantuml"
- C4-PlantUML (text C4 syntax)  → diagram_type="c4plantuml"
- GraphViz dot graphs           → diagram_type="graphviz"
- BPMN process flows            → diagram_type="bpmn"
The tool returns a markdown image tag — include it directly in your response.

── POST-TOOL RESPONSE ────────────────────────────────────────────────────────
After a generate_architecture or generate_architecture_from_context tool call
completes, the full document has already been streamed to the user in real time.
Respond with 1-2 sentences ONLY — e.g. "Architecture document generated. Would
you like me to save it as a .docx file?" Do NOT repeat or summarise the tool
output.

── MERMAID THAT ACTUALLY RENDERS ─────────────────────────────────────────────
A diagram that fails to parse shows the reader an error box instead of a diagram,
so it is worse than a simpler diagram that works.

USE ONLY these four types — they are the ones the templates use and the ones the
viewer renders reliably:
    graph / flowchart      classDiagram      sequenceDiagram      erDiagram

NEVER use C4Context, C4Container, C4Component, mindmap, timeline, quadrantChart or
block-beta. Mermaid's C4 support is experimental and fails often; the C4 sections
above are drawn with `graph TB` for exactly that reason — follow those templates
rather than reaching for a C4-specific syntax.

Rules that prevent most parse failures:
- Put EVERY node label in double quotes: A["Order Service"], never A[Order Service].
  An unquoted label containing a bracket, parenthesis, comma, colon or slash is the
  single most common cause of a broken diagram.
- Never put a raw " inside a label. Use a different word.
- In classDiagram, keep method signatures simple: +save(), +findById(). Avoid
  generics (List~T~), nested parentheses and return-type annotations.
- One statement per line. Do not put a comment on the same line as a statement.
- Keep every diagram under ~25 nodes; split it rather than growing it.

── ABSOLUTE RULES ────────────────────────────────────────────────────────────
- NEVER write architecture as a plain-text reply. ALWAYS call a tool first.
- NEVER say you cannot access files — always call read_document.
- EVERY architecture output MUST have at least one ```mermaid``` block.
- HIGH-LEVEL DESIGN MUST include a high-level architecture diagram (Mermaid graph/flowchart).
- LOW-LEVEL DESIGN MUST include BOTH a component/class diagram (Mermaid classDiagram
  or detailed flowchart) AND a sequence diagram (Mermaid sequenceDiagram) — never omit either.
- C4 MUST show all three levels (Context, Container, Component).
- DB schema MUST include both an ER diagram (Mermaid erDiagram) AND CREATE TABLE SQL.
- API CONTRACT MUST include full OpenAPI 3.0 YAML.
- EVERY architecture output MUST include the ## SECURITY DESIGN CHECKLIST section
  with all 10 OWASP categories assessed against THIS design.
- Use descriptive filenames when saving: 'leave_mgmt_hld.docx'.
- NEVER invent API endpoints, database tables, business rules, or user roles
  that are not present in the requirements. Tag any necessary assumption with
  ⚠️ [ASSUMPTION] so the user can confirm before development starts.
- IF requirements context is empty or has fewer than 3 user stories/features,
  ask the user for more detail before generating — do not produce a generic
  architecture diagram that doesn't match their actual system.

── ADO WRITE-BACK ─────────────────────────────────────────────────────────────
After the user confirms the design is complete AND a DOCX has been saved:
1. Call update_ado_epic_design_complete(project, epic_id, design_url, target_state)
   - project: ADO project name (ask the user if not known)
   - epic_id: work item ID of the parent epic (ask the user if not known)
   - design_url: the URL returned by save_architecture
   - target_state: "Design Complete" (default) — confirm real state name if unsure
2. The tool moves the epic state and adds a comment with the document link.
3. Always confirm with the user before calling this tool.
4. If the user hasn't saved a document yet, call save_architecture first.

── PIPELINE HANDOFF ─────────────────────────────────────────────────────────
Trigger: "start coding", "develop it", "write the code", "build it",
"proceed to development", "next step" (after design is confirmed complete).
Response: one sentence confirming design is complete, then call write_design_artifact
with the run_id from the session context and ALL relevant sections populated,
INCLUDING security_checklist (the OWASP Top-10 review).

── FINAL REMINDER (HIGHEST PRIORITY) ─────────────────────────────────────────
The user has already provided requirements/user stories in the structured pipeline
context. When asked to generate/design/save, follow this EXACT tool sequence:
  STEP 1 — call `generate_architecture_from_context` with the user stories passed as
           the `context` argument (a single string). This returns the full markdown
           architecture document. Do NOT call markdowntodoc or save_architecture yet.
  STEP 2 — take the markdown the tool returned and call `save_architecture` with that
           markdown as `content` and a `filename` like "design.docx".
NEVER call `markdowntodoc` directly, and NEVER call any save tool with an empty
`content` — you must generate the architecture in STEP 1 first. NEVER end a turn with
only "I'll generate it now" and no tool call: that is a hard failure. Do not ask for
requirements you already have.
"""

# Append the shared, agent-agnostic MCP provenance note (see shared/tools/mcp_runtime).
DESIGN_SYS_MESSAGE = DESIGN_SYS_MESSAGE + MCP_TOOLS_PROMPT_NOTE


# ── LangGraph workflow ─────────────────────────────────────────────────────────

class AgentState(TypedDict):
    messages: Annotated[List[BaseMessage], add_messages]
    tenant_id: str
    model_id: str | None
    offering_id: str | None
    # BYOK-resolved model carried through state so the tools node can re-establish it
    # (LangGraph runs nodes in separate task contexts, so the contextvar set in the
    # agent node does not reliably reach the tools that call _llm_generate_async).
    resolved_model: Any


# Per-alias orchestrator cache — keyed by (alias, model) to avoid re-constructing
# ChatLiteLLM on every invocation while still supporting per-tenant BYOK keys (Pitfall 4).
_ORCHESTRATOR_CACHE: dict[tuple[str, str], object] = {}


def _build_orchestrator(model: str, litellm_provider: str, api_key: str,
                        base_url: str | None, alias: str) -> object:
    """Build (or return cached) ChatLiteLLM orchestrator. Direct-provider BYOK call.
    Cached by (alias, model); alias is non-secret (SC#5)."""
    cache_key = (alias, model)
    if cache_key in _ORCHESTRATOR_CACHE:
        return _ORCHESTRATOR_CACHE[cache_key]
    # Deferred: importing litellm costs ~7s. sys.modules makes repeat calls free.
    from langchain_litellm import ChatLiteLLM
    from shared.services.model_resolver import temperature_kwargs  # noqa: PLC0415

    instance = ChatLiteLLM(
        model=model,
        custom_llm_provider=litellm_provider,
        api_base=base_url,
        api_key=api_key,
        max_tokens=8192,
        # Omitted entirely for gpt-5-family models, which reject any temperature but
        # their default — see temperature_kwargs.
        **temperature_kwargs(model, 0.2),
        # 0, not langchain_litellm's own default -- guarded_completion
        # (shared/services/model_call_wrapper.py) already retries the whole
        # ainvoke call; ChatLiteLLM's own tenacity retry underneath it also
        # retries on RateLimitError, uncoordinated with the outer loop. See
        # dev_agent.py::_build_llm's identical fix and "desicions and
        # issues.txt" Issue 14 for the live evidence.
        max_retries=0,
        # Stream tokens so the copilot shows the design agent's replies live.
        streaming=True,
    )
    _ORCHESTRATOR_CACHE[cache_key] = instance
    return instance


def _sanitize_messages(messages: list) -> list:
    """Keep tool_use / tool_result blocks paired in BOTH directions.

    The implementation moved to `shared.services.message_pairing` unchanged — the
    Development and Requirements agents need exactly this repair for exactly the same
    reason, and three copies is how one of them ends up with only half the rule (which
    is what had already happened). See that module for why the invariant matters and
    why a dangling call is stripped rather than answered with a fabricated result.
    """
    from shared.services.message_pairing import sanitize_tool_call_pairing  # noqa: PLC0415

    return sanitize_tool_call_pairing(messages)


async def agent(state: AgentState):
    from shared.services.model_resolver import resolve_model_for_run, set_resolved_model, NoModelConfiguredError, ModelNotEnabledError
    tenant_id = state.get("tenant_id", "")
    requested = state.get("model_id")
    try:
        resolved = await resolve_model_for_run(tenant_id, requested, offering_id=state.get("offering_id"))
    except (NoModelConfiguredError, ModelNotEnabledError) as e:
        logger.warning("Design agent model resolution failed (tenant=%s): %s", tenant_id, type(e).__name__)
        return {"messages": [AIMessage(content=(
            "No usable model is configured for your organization. "
            "An administrator must add and verify a model provider in Org Settings → Model Providers."))]}
    except Exception as e:
        logger.error("Design agent model resolution error (tenant=%s): %s", tenant_id, type(e).__name__)
        # The TYPE NAME ALONE WAS THE WHOLE REPLY — a user hit an Azure rate limit
        # and was shown "Agent error: RateLimitError", which does not say it is
        # temporary, not their fault, or that another model would work now.
        # Still never str(exc): a BYOK provider error can echo the tenant key.
        from shared.services.model_errors import friendly_model_error  # noqa: PLC0415

        return {"messages": [AIMessage(content=friendly_model_error(e))]}
    set_resolved_model(resolved)
    try:
        clean_messages = _sanitize_messages(state["messages"])
        orch = _build_orchestrator(resolved.model, resolved.litellm_provider,
                                   resolved.api_key, resolved.base_url, resolved.alias)
        # Bind base tools + any per-run MCP tools here (not in the cached builder) so
        # run-specific MCP tools never leak across runs via the shared orchestrator cache.
        orch = orch.bind_tools(tools + get_skill_tools("design") + get_mcp_tools())
        from shared.services.model_call_wrapper import guarded_completion

        response = await guarded_completion(
            resolved, orch, clean_messages, tenant_id=tenant_id, agent_type="design",
            config={"metadata": {"user_api_key_alias": resolved.alias}},
        )
        return {"messages": [response], "resolved_model": resolved}
    except Exception as e:
        # NEVER str(exc) here — credential leakage risk (BYOK-resolved model calls
        # can surface the tenant's own API key in a provider error message).
        logger.error(
            "Design agent error (tenant=%s alias=%s): %s",
            tenant_id,
            resolved.alias,
            type(e).__name__,
        )
        # The TYPE NAME ALONE WAS THE WHOLE REPLY — a user hit an Azure rate limit
        # and was shown "Agent error: RateLimitError", which does not say it is
        # temporary, not their fault, or that another model would work now.
        # Still never str(exc): a BYOK provider error can echo the tenant key.
        from shared.services.model_errors import friendly_model_error  # noqa: PLC0415

        return {"messages": [AIMessage(content=friendly_model_error(e))]}


_TOOL_ARG_HINTS: dict[str, str] = {
    "markdowntodoc": (
        "Do NOT call markdowntodoc directly. First call generate_architecture_from_context "
        "(pass the user stories as 'context') to produce the architecture markdown, then "
        "call save_architecture with that markdown as 'content' and a 'filename'."
    ),
    "save_architecture": (
        "save_architecture requires two arguments: "
        "'content' (the full markdown text to save) and "
        "'filename' (e.g. 'architecture.docx'). "
        "Pass the architecture markdown you just generated as 'content'."
    ),
    "generate_architecture": (
        "generate_architecture requires 'requirements_text' (the extracted requirements text). "
        "Read the source document first with read_document, then pass the result here."
    ),
    "render_diagram_via_kroki": (
        "render_diagram_via_kroki requires 'diagram_type' (e.g. 'plantuml') and 'source' (the raw diagram source code)."
    ),
}


async def action(state: AgentState):
    # Re-establish the BYOK-resolved model in this node's context — LangGraph may run
    # the tools node in a separate task context where the agent node's contextvar set
    # didn't propagate, which would make tools calling _llm_generate_async fail
    # "No model resolved". Read the already-resolved model from state (no DB call).
    try:
        from shared.services.model_resolver import set_resolved_model, get_resolved_model
        if get_resolved_model() is None and state.get("resolved_model") is not None:
            set_resolved_model(state["resolved_model"])
    except Exception:  # noqa: BLE001 — defensive; tools degrade with a clear message
        pass
    last_message = state["messages"][-1]
    results = []
    for tc in getattr(last_message, "tool_calls", []):
        broadcast_log(manager, f"Invoking tool {tc['name']} with args {tc['args']}", level="LOGS")
        try:
            tool_to_call = next(
                (t for t in tools + get_skill_tools("design") + get_mcp_tools()
                 if t.name == tc["name"]), None
            )
            if tool_to_call is None:
                obs = f"Unknown tool: {tc['name']}"
            else:
                obs = await tool_to_call.ainvoke(tc["args"])
        except Exception as e:
            import traceback  # noqa: PLC0415
            hint = _TOOL_ARG_HINTS.get(tc["name"], "")
            # Log the REAL error (type + traceback) server-side; keep the chat message
            # concise so a huge exception string can't flood the UI as "content".
            # traceback.format_exc() already carries the exception's own message as
            # its last line — passing e too would be the same str(exc) exposure the
            # rest of this codebase avoids (credential leakage risk).
            logger.error(
                "design tool %s failed: %s\n%s",
                tc["name"], type(e).__name__, traceback.format_exc(),
            )
            short = str(e)
            if len(short) > 400:
                short = short[:400] + "… (truncated — see server logs)"
            obs = f"Tool '{tc['name']}' failed: {type(e).__name__}: {short}" + (f" — {hint}" if hint else "")
        results.append(ToolMessage(content=str(obs), tool_call_id=tc["id"], name=tc["name"]))
    return {"messages": results}


def should_continue(state: AgentState) -> str:
    last = state["messages"][-1]
    return "action" if getattr(last, "tool_calls", None) else False


workflow = StateGraph(AgentState)
workflow.add_node("agent", agent)
workflow.add_node("tools", action)
workflow.add_edge("tools", "agent")
workflow.set_entry_point("agent")
workflow.add_conditional_edges("agent", should_continue, {"action": "tools", False: END})

app = workflow.compile(checkpointer=_build_checkpointer("design"))


async def run_architecture_workflow_async(messages, config=None):
    if config is None:
        config = {"configurable": {"thread_id": "default"}}
    result = None
    async for event in app.astream({"messages": messages}, config):
        result = event
    return result


def run_architecture_workflow_sync(messages, config=None):
    return asyncio.run(run_architecture_workflow_async(messages, config))
